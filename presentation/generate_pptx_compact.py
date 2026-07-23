#!/usr/bin/env python3
"""
AETHERIX Compact Animated Presentation Generator
Creates a 25-slide PPTX with only the compact slide set.
"""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
import copy
from lxml import etree

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_DIR = os.path.join(BASE_DIR, "presentation", "output")
CHARTS_DIR = os.path.join(BASE_DIR, "visualizations", "charts")
DIAGRAMS_DIR = os.path.join(BASE_DIR, "visualizations", "diagrams")
IMG_DIR = os.path.join(BASE_DIR, "docs", "img")

os.makedirs(OUTPUT_DIR, exist_ok=True)

BG_DARK = RGBColor(0x0B, 0x0E, 0x1A)
BG_SLIDE = RGBColor(0x0D, 0x12, 0x22)
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xAA)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x4D, 0x4D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xCC)
MED_GRAY = RGBColor(0x6B, 0x7B, 0x96)
CARD_BG = RGBColor(0x14, 0x1B, 0x2D)
CARD_BORDER = RGBColor(0x1E, 0x2A, 0x42)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
YELLOW = RGBColor(0xFF, 0xD9, 0x3D)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

TOTAL_SLIDES = 31

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text="", font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_card(slide, left, top, width, height, title="", body_lines=None, title_color=ACCENT_BLUE, body_color=LIGHT_GRAY, title_size=16, body_size=13, fill=CARD_BG, border=CARD_BORDER):
    shape = add_shape(slide, left, top, width, height, fill_color=fill, border_color=border)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Calibri"
    if body_lines:
        for line in body_lines:
            add_paragraph(tf, line, font_size=body_size, color=body_color, space_before=Pt(2), space_after=Pt(2))
    return shape


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    return add_shape(slide, left, top, width, height, fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    if os.path.exists(img_path):
        if width and height:
            return slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            return slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            return slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            return slide.shapes.add_picture(img_path, left, top)
    return None


def add_table(slide, left, top, width, height, rows, cols, data, header_color=ACCENT_BLUE, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = data[row_idx][col_idx] if row_idx < len(data) and col_idx < len(data[row_idx]) else ""
            cell.fill.solid()
            if row_idx == 0:
                cell.fill.fore_color.rgb = header_color
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                    p.font.size = Pt(12)
                    p.font.name = "Calibri"
                    p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 1 else RGBColor(0x18, 0x22, 0x36)
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = LIGHT_GRAY
                    p.font.size = Pt(11)
                    p.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    return table_shape


def add_entrance_animation(slide, shape, delay_ms=0, duration_ms=500, anim_type="fade"):
    slide_element = slide._element
    timing = slide_element.find(qn('p:timing'))
    if timing is None:
        timing = parse_xml(
            '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            '<p:childTnLst/>'
            '</p:cTn></p:par></p:tnLst>'
            '</p:timing>'
        )
        slide_element.append(timing)
    child_tn_lst = timing.find('.//' + qn('p:childTnLst'))
    shape_id = shape.shape_id
    seq_id = len(child_tn_lst) + 2
    if anim_type == "fade":
        anim_xml = f'<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><p:cTn id="{seq_id}" fill="hold"><p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{seq_id+1}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{seq_id+2}" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" nodeType="afterEffect"><p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id="{seq_id+3}" dur="{duration_ms}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set><p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{seq_id+4}" dur="{duration_ms}"/><p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl></p:cBhvr></p:animEffect></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
    else:
        anim_xml = f'<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="{seq_id}" fill="hold"><p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{seq_id+1}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{seq_id+2}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect"><p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id="{seq_id+3}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
    child_tn_lst.append(parse_xml(anim_xml))


def add_slide_transition(slide, transition_type="fade", duration_ms=700):
    slide_element = slide._element
    trans = {
        "fade": '<p:fade/>',
        "push": '<p:push dir="l"/>',
        "wipe": '<p:wipe dir="d"/>',
        "cover": '<p:cover dir="l"/>',
    }
    inner = trans.get(transition_type, '<p:fade/>')
    trans_xml = f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advClick="1">{inner}</p:transition>'
    existing = slide_element.find(qn('p:transition'))
    if existing is not None:
        slide_element.remove(existing)
    slide_element.insert(1, parse_xml(trans_xml))


_footer_counter = [1]


def add_footer(slide, slide_num=None, total=None, citations=None):
    _footer_counter[0] += 1
    n = _footer_counter[0]
    if citations:
        add_textbox(slide, Inches(0.5), Inches(7.22), Inches(12.5), Inches(0.25),
                    citations, font_size=8, color=MED_GRAY)
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(5), Inches(0.4),
                "AETHERIX \u2014 Interplanetary Communication Network", font_size=10, color=MED_GRAY)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4),
                f"{n} / {TOTAL_SLIDES}", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


def add_stat_card(slide, x, y, w, h, value, label, value_color=ACCENT_BLUE, val_size=24):
    card = add_shape(slide, x, y, w, h, fill_color=CARD_BG, border_color=value_color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(val_size)
    p.font.color.rgb = value_color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(tf, label, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(2))
    return card


def new_slide():
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    return slide


def add_chart_slide(chart_file, title, subtitle, caption, accent_rgb, citations=None):
    slide = new_slide()
    add_slide_transition(slide, "fade")
    add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7),
                title, font_size=32, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.7), Inches(0.9), Inches(11.0), Inches(0.4),
                subtitle, font_size=16, color=accent_rgb)
    add_accent_line(slide, Inches(0.7), Inches(1.35), Inches(2.5), accent_rgb)
    add_image_safe(slide, chart_file, Inches(0.7), Inches(1.5), Inches(7.0), Inches(4.5))
    add_textbox(slide, Inches(0.7), Inches(6.2), Inches(11.0), Inches(0.3),
                caption, font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    add_footer(slide, citations=citations)
    return slide


_SPEAKER_NOTES = {
    1: "State your name clearly. Read the topic number and title exactly as on the exam paper. Pause to let examiners see it. (30 seconds)",
    2: "Quick overview of what we will cover. 12 topics across 25 slides. (20 seconds)",
    3: "This slide sets up the narrative arc. First explain what AETHERIX is in plain language - it's like the postal service for interplanetary space. (1.5 minutes)",
    4: "Start with the scale. 54.6M to 401M km. Light itself takes 3-22 minutes one way. TCP/IP was designed for sub-second round trips. (1.5 minutes)",
    5: "The key insight: instead of requiring an end-to-end connection like TCP, DTN works like the postal service. Each node takes custody of your data. (1.5 minutes)",
    6: "Show the architecture. Six core modules feed into the simulation engine, which feeds the web showcase. Point to each module as you explain. (1 minute)",
    7: "Architecture diagram showing source modules feeding simulation engine and web demos.",
    8: "BPv7 is the postal service protocol. Walk through the stack: Application at top, BPv7 as the store-and-forward layer, three convergence layers. (2 minutes)",
    9: "241 nodes across 5 tiers. Walk through each tier. Earth Ground is the DSN. Earth Orbital has LEO laser mesh. Deep Space has Lagrange point relays. (2 minutes)",
    10: "Visual overview of the 5-tier topology with 3 redundant paths.",
    11: "RUN THE LIVE DEMO from the Link Budget page. Show the 3 distance scenarios. 1550nm was chosen for telecom heritage. (2 minutes)",
    12: "Data rate degrades with distance squared. (15 seconds)",
    13: "Walk through the 7-hop journey. 500MB from Perseverance to JPL. Total transit ~13 min vs 12.5 min light-time. (2 minutes)",
    14: "CGR is what NASA uses today. It's static. Our RL agent learns from experience. Multi-agent federated learning means agents at each node share knowledge. (2 minutes)",
    15: "RL agent Q-value heatmap showing learned routing preferences. (15 seconds)",
    16: "BB84 is beautifully simple: send qubits, measure, compare bases, check QBER. If QBER is below 11%, no one listened in. (2 minutes)",
    17: "QKD security analysis with QBER threshold. (15 seconds)",
    18: "Mars and Earth dance around the Sun with a 26-month synodic period. Everything changes - distance, delay, bandwidth. (1.5 minutes)",
    19: "Visual data flow through the protocol stack.",
    20: "Hit these numbers with confidence. 10-100x faster. >95% availability vs 60-75%. Quantum-secure. (1 minute)",
    21: "Side-by-side performance comparison chart. (15 seconds)",
    22: "This is real, working code. 27 Python modules, 480 tests, 12 interactive demos. All the physics is real. (1.5 minutes)",
    23: "Phases 1-4 are done. Phase 5: ns-3 simulation. Phase 6: Upgrade to DQN and integrate with ION-DTN. (1.5 minutes)",
    24: "Summarize the problem and solution clearly. Point to the numbers. Offer to show live demos. (1 minute)",
    25: "Summarize the four key numbers: 10-100x faster, >95% availability, AI-powered routing, quantum-secure. Invite questions. (30 seconds)",
}


# ============================================================
# SLIDE 1 — Introduction (Title Hero)
# ============================================================
print("Creating Slide 1: Introduction...")
slide = new_slide()
add_slide_transition(slide, "fade")
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.RECTANGLE)
add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.5), "AETHERIX", font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(4.5), Inches(2.5), Inches(4.3), ACCENT_CYAN, Pt(4))
add_textbox(slide, Inches(1.0), Inches(2.8), Inches(11.3), Inches(1.0), "Autonomous Extraterrestrial High-throughput Enhancing Routing\nand Inter-planetary eXchange", font_size=22, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.0), Inches(4.0), Inches(9.3), Inches(0.5), "EduQual Level 6 Diploma in AI Operations  |  Topic 59", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.0), Inches(4.5), Inches(9.3), Inches(0.5), "Building an Interplanetary Communication Network with DTN,\nQuantum Communication, and Space-Based Infrastructure for Mars Mission Support", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(3.5), Inches(5.3), Inches(6.3), ACCENT_BLUE, Pt(2))
add_textbox(slide, Inches(2.0), Inches(5.6), Inches(9.3), Inches(0.5), "Muhammad Abdullah Tariq", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.0), Inches(6.1), Inches(9.3), Inches(0.5), "September 2026", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.0), Inches(6.7), Inches(9.3), Inches(0.4), "matx104.github.io/AETHERIX  |  github.com/matx104/AETHERIX", font_size=12, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.RECTANGLE)


# ============================================================
# SLIDE 2 — Agenda
# ============================================================
print("Creating Slide 2: Agenda...")
slide = new_slide()
add_slide_transition(slide, "push")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "PRESENTATION AGENDA", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_footer(slide, 2, citations="[A2] AETHERIX topology.py (241 nodes)  \u00b7  References slides at end of deck")

agenda_items = [
    ("01", "The Challenge", "Why space breaks the internet", ACCENT_BLUE),
    ("02", "Architecture", "DTN + AI + Quantum Security", ACCENT_CYAN),
    ("03", "DTN & BPv7", "Store-and-forward foundation", ACCENT_PURPLE),
    ("04", "Topology", "241 nodes across two worlds", ACCENT_BLUE),
    ("05", "Link Budget", "1550nm laser analysis", ACCENT_ORANGE),
    ("06", "RL Routing", "Multi-agent federated Q-learning", ACCENT_CYAN),
    ("07", "Quantum Security", "BB84/E91 + repeater chains", ACCENT_PURPLE),
    ("08", "Orbital Mechanics", "Contact windows & synodic period", ACCENT_BLUE),
    ("09", "Mars Mission", "End-to-end simulation walkthrough", ACCENT_ORANGE),
    ("10", "Performance", "AETHERIX vs current systems", GREEN),
    ("11", "Roadmap", "CCSDS, IETF, deployment phases", ACCENT_CYAN),
    ("12", "Q&A", "Summary and live demo", WHITE),
]

_per_col = 6
for i, (num, title, desc, color) in enumerate(agenda_items):
    col_x = Inches(0.7) if i < _per_col else Inches(6.8)
    row_y = Inches(1.25) + Inches(0.80) * (i % _per_col)
    card = add_card(slide, col_x, row_y, Inches(5.8), Inches(0.72), border=color)
    add_textbox(slide, col_x + Inches(0.15), row_y + Inches(0.06), Inches(5.5), Inches(0.32), f"{num}  {title}", font_size=13, color=WHITE, bold=True)
    add_textbox(slide, col_x + Inches(0.15), row_y + Inches(0.38), Inches(5.5), Inches(0.3), desc, font_size=10, color=LIGHT_GRAY)
    add_entrance_animation(slide, card, delay_ms=80 * i, anim_type="fade")


# ============================================================
# SLIDE 3 — What is AETHERIX
# ============================================================
print("Creating Slide 3: What is AETHERIX...")
slide = new_slide()
add_slide_transition(slide, "fade")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "WHAT IS AETHERIX?", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11.0), Inches(0.4), "Overview & The Problem", font_size=16, color=MED_GRAY)
add_footer(slide, 3, citations="[1] NASA MRN 2024  \u00b7  [3] JPL Horizons  \u00b7  [12] RFC 4838  \u00b7  [A2] topology.py")

card_left = add_card(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(3.8), border=ACCENT_BLUE)
add_textbox(slide, Inches(0.9), Inches(1.7), Inches(5.2), Inches(0.45), "What is AETHERIX?", font_size=18, color=ACCENT_BLUE, bold=True)
add_accent_line(slide, Inches(0.9), Inches(2.1), Inches(2.0), ACCENT_BLUE, Pt(2))
what_lines = [
    "Autonomous Extraterrestrial High-throughput Enhancing\nRouting and Inter-planetary eXchange",
    "An architecture for delay-tolerant networking (DTN)\nbetween Earth and Mars",
    "Replaces static Contact Graph Routing with\nreinforcement-learning-based adaptive routing",
    "Secures links via Quantum Key Distribution\n(BB84 & E91 protocols with repeater chains)",
    "Hybrid optical/RF communications across\na 5-tier, 241-node interplanetary topology",
    "Open-source, standards-compliant (CCSDS, IETF)\nresearch platform for deep-space networking",
]
for j, line in enumerate(what_lines):
    add_textbox(slide, Inches(0.9), Inches(2.2) + Inches(0.5) * j, Inches(5.2), Inches(0.5), line, font_size=10, color=LIGHT_GRAY)

card_right = add_card(slide, Inches(6.7), Inches(1.6), Inches(5.6), Inches(3.8), border=ACCENT_RED)
add_textbox(slide, Inches(6.9), Inches(1.7), Inches(5.2), Inches(0.45), "The Problem", font_size=18, color=ACCENT_RED, bold=True)
add_accent_line(slide, Inches(6.9), Inches(2.1), Inches(2.0), ACCENT_RED, Pt(2))
problem_lines = [
    "Earth\u2013Mars distance varies from 54.6M to 401M km\n\u2014 a 7x range swing every 780 days",
    "One-way light delay is 3\u201322 minutes;\nround-trip TCP ACK is 6\u201344 minutes",
    "Solar conjunction causes ~2-week communication\nblackouts twice per synodic period",
    "Deep-space links are inherently intermittent:\nno continuous end-to-end path exists",
    "Static routing cannot adapt to dynamic\ncontact schedules and link degradation",
    "Classical key exchange is infeasible across\ninterplanetary distances",
]
for j, line in enumerate(problem_lines):
    add_textbox(slide, Inches(6.9), Inches(2.2) + Inches(0.5) * j, Inches(5.2), Inches(0.5), line, font_size=10, color=LIGHT_GRAY)

callout_shape = add_shape(slide, Inches(0.7), Inches(5.7), Inches(11.6), Inches(0.9), fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(1.0), Inches(5.8), Inches(11.0), Inches(0.7), "AETHERIX addresses every one of these challenges \u2014 harnessing DTN store-and-forward, AI-driven adaptive routing, quantum-secured key distribution, and hybrid optical/RF links.", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4 — The Distance
# ============================================================
print("Creating Slide 4: The Distance...")
slide = new_slide()
add_slide_transition(slide, "push")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "THE DISTANCE", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_RED, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11.0), Inches(0.4), "Why Space Breaks the Internet", font_size=16, color=MED_GRAY)
add_footer(slide, 4, citations="[3] JPL Horizons (distance/light-time)  \u00b7  [12] RFC 4838 (DTN rationale)  \u00b7  [1] NASA MRO data rate")

left_headers = ["Parameter", "Value"]
left_rows = [
    ["Closest approach", "54.6 million km"],
    ["Farthest distance", "401 million km"],
    ["Distance variation", "7.3\u00d7 (735%)"],
    ["One-way light time", "3 \u2013 22 minutes"],
    ["Synodic period", "780 days (26 months)"],
    ["Solar conjunction blackout", "~2 weeks (twice/period)"],
]
_d1 = [left_headers] + left_rows
add_table(slide, Inches(0.7), Inches(1.5), Inches(5.6), Inches(3.2), len(_d1), len(left_headers), _d1, header_color=ACCENT_RED)

right_headers = ["TCP/IP Assumption", "Deep-Space Reality"]
right_rows = [
    ["Low latency (< 200 ms)", "6 \u2013 44 min round trip"],
    ["Continuous connectivity", "Scheduled contact windows only"],
    ["End-to-end path exists", "No simultaneous path possible"],
    ["Reliable ACKs in seconds", "ACKs take minutes to hours"],
]
_d2 = [right_headers] + right_rows
add_table(slide, Inches(6.7), Inches(1.5), Inches(5.6), Inches(2.2), len(_d2), len(right_headers), _d2, header_color=ACCENT_ORANGE)

quote_card = add_card(slide, Inches(6.7), Inches(4.0), Inches(5.6), Inches(1.3), border=ACCENT_ORANGE)
add_textbox(slide, Inches(6.9), Inches(4.1), Inches(5.2), Inches(1.1), '"TCP/IP assumes the network is fast, reliable, and always connected.\nDeep space is none of those things."', font_size=12, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

stat_y = Inches(5.6)
stat_labels = ["54.6M km", "401M km", "780 days", "0.5\u20136 Mbps"]
stat_sublabels = ["Min Distance", "Max Distance", "Synodic Period", "Current DSN Rate"]
stat_colors = [ACCENT_BLUE, ACCENT_RED, ACCENT_ORANGE, ACCENT_CYAN]
for k in range(4):
    sx = Inches(0.7) + Inches(3.0) * k
    add_shape(slide, sx, stat_y, Inches(2.7), Inches(1.1), fill_color=CARD_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_accent_line(slide, sx, stat_y, Inches(2.7), stat_colors[k], Pt(3))
    add_textbox(slide, sx + Inches(0.15), stat_y + Inches(0.15), Inches(2.4), Inches(0.45), stat_labels[k], font_size=22, color=stat_colors[k], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, sx + Inches(0.15), stat_y + Inches(0.6), Inches(2.4), Inches(0.35), stat_sublabels[k], font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — The Answer
# ============================================================
print("Creating Slide 5: The Answer...")
slide = new_slide()
add_slide_transition(slide, "wipe")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "THE ANSWER", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11.0), Inches(0.4), "Delay-Tolerant Networking", font_size=16, color=MED_GRAY)
add_footer(slide, 5, citations="[9] RFC 9171 BPv7  \u00b7  [10] RFC 5326 LTP  \u00b7  [11] RFC 7242 TCPCL  \u00b7  [12] RFC 4838")

flow_labels = ["Bundle", "Store", "Wait", "Forward", "Deliver"]
flow_colors = [ACCENT_BLUE, ACCENT_PURPLE, ACCENT_ORANGE, ACCENT_CYAN, GREEN]
flow_descs = [
    "Create data\nbundle (BPv7)",
    "Buffer at\nlocal node",
    "Hold until\ncontact window",
    "Transmit to\nnext hop",
    "Reach final\ndestination",
]
card_w = Inches(2.0)
card_h = Inches(1.6)
start_x = Inches(0.5)
gap = Inches(0.35)
arrow_w = Inches(0.3)
flow_y = Inches(1.6)

for idx, (label, color, desc) in enumerate(zip(flow_labels, flow_colors, flow_descs)):
    cx = start_x + (card_w + arrow_w + gap) * idx
    card = add_card(slide, cx, flow_y, card_w, card_h, border=color)
    add_textbox(slide, cx + Inches(0.1), flow_y + Inches(0.1), card_w - Inches(0.2), Inches(0.35), label, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, cx + Inches(0.1), flow_y + Inches(0.5), card_w - Inches(0.2), Inches(0.9), desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_entrance_animation(slide, card, delay_ms=150 * idx, anim_type="fade")
    if idx < len(flow_labels) - 1:
        arrow_x = cx + card_w + Inches(0.05)
        add_textbox(slide, arrow_x, flow_y + Inches(0.5), arrow_w, Inches(0.5), "\u2192", font_size=28, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

stat_card_data = [
    ("BPv7", "Bundle Protocol v7\nRFC 9171 \u2014 store &\nforward with custody", ACCENT_BLUE),
    ("AI Routing", "Multi-agent federated\nQ-learning adapts in\nreal time to contacts", ACCENT_PURPLE),
    ("QKD", "BB84 & E91 quantum\nkey distribution with\nrepeater chains", ACCENT_CYAN),
]
stat_y = Inches(3.5)
for idx, (title, desc, color) in enumerate(stat_card_data):
    sx = Inches(0.7) + Inches(4.0) * idx
    add_stat_card(slide, sx, stat_y, Inches(3.6), Inches(1.4), title, desc, value_color=color)

postal_card = add_card(slide, Inches(0.7), Inches(5.2), Inches(11.6), Inches(1.0), border=ACCENT_BLUE)
add_textbox(slide, Inches(1.0), Inches(5.3), Inches(11.0), Inches(0.8), 'Think of it like the postal service: you don\'t need a continuous connection between sender and receiver \u2014 each post office (node) stores the letter (bundle) until the next truck (contact window) is available, then forwards it one hop closer to the destination.', font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6 — System Architecture
# ============================================================
print("Creating Slide 6: System Architecture...")
slide = new_slide()
add_slide_transition(slide, "cover")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "SYSTEM ARCHITECTURE", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11.0), Inches(0.4), "Five Core Modules", font_size=16, color=MED_GRAY)
add_footer(slide, 6, citations="[A1] rl_agent.py  \u00b7  [A2] topology.py  \u00b7  [A3] simulator.py  \u00b7  [A4] link_budget.py  \u00b7  [A5] qkd.py")

modules = [
    ("1", "Infrastructure", "Link budget & RF/optical analysis", ACCENT_BLUE),
    ("2", "Routing", "RL agent, BPv7, store-and-forward", ACCENT_PURPLE),
    ("3", "Security", "QKD (BB84/E91) & repeater chains", ACCENT_CYAN),
    ("4", "Orbital", "Contact windows & Doppler analysis", ACCENT_ORANGE),
    ("5", "Simulation", "Full network simulation engine", GREEN),
]

for idx, (num, name, desc, color) in enumerate(modules):
    my = Inches(1.6) + Inches(0.95) * idx
    circle = add_shape(slide, Inches(0.7), my, Inches(0.6), Inches(0.6), fill_color=color, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide, Inches(0.7), my + Inches(0.08), Inches(0.6), Inches(0.45), num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), my + Inches(0.0), Inches(4.0), Inches(0.35), name, font_size=16, color=color, bold=True)
    add_textbox(slide, Inches(1.5), my + Inches(0.35), Inches(4.0), Inches(0.35), desc, font_size=11, color=LIGHT_GRAY)

table_headers = ["Module", "Engine", "Showcase"]
table_rows = [
    ["Infrastructure", "OpticalLinkBudget + RFLinkBudget", "Earth-Mars link analysis"],
    ["Routing", "RLRoutingAgent + ForwardingEngine", "Multi-hop DTN delivery"],
    ["Security", "BB84Protocol + E91Protocol", "Quantum-secured key exchange"],
    ["Orbital", "ContactWindows + Doppler", "Synodic contact prediction"],
    ["Simulation", "Simulator + PolicyEngine", "End-to-end scenario runs"],
]
_d4 = [table_headers] + table_rows
add_table(slide, Inches(6.4), Inches(1.6), Inches(6.0), Inches(4.5), len(_d4), len(table_headers), _d4, header_color=ACCENT_BLUE)


# ============================================================
# SLIDE 7 — Architecture Diagram
# ============================================================
print("Creating Slide 7: Architecture Diagram...")
slide = new_slide()
add_slide_transition(slide, "fade")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "ARCHITECTURE DIAGRAM", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_footer(slide, 7, citations="[A2] AETHERIX topology.py (5-tier, 241 nodes)  \u00b7  github.com/matx104/AETHERIX")

arch_img_path = os.path.join(DIAGRAMS_DIR, "system_architecture.png")
add_image_safe(slide, arch_img_path, Inches(0.7), Inches(1.2), Inches(11.6), Inches(5.8))


# ============================================================
# SLIDE 8 — BPv7 Deep Dive
# ============================================================
print("Creating Slide 8: BPv7 Deep Dive...")
slide = new_slide()
add_slide_transition(slide, "push")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "BPv7 DEEP DIVE", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11.0), Inches(0.4), "The Foundation", font_size=16, color=MED_GRAY)
add_footer(slide, 8, citations="[9] RFC 9171  \u00b7  [2] CCSDS 734.2-B-1  \u00b7  [10] RFC 5326 LTP  \u00b7  [11] RFC 7242 TCPCL  \u00b7  [12] RFC 4838")

stack_layers = [
    ("Application Layer", "User data, commands, telemetry", ACCENT_BLUE),
    ("Bundle Protocol v7", "Routing, custody, fragmentation, TTL", ACCENT_PURPLE),
    ("Convergence Layer", "LTP (space), TCPCL (Earth), UDPCL (ISL)", ACCENT_CYAN),
    ("Physical Layer", "Optical 1550nm / Ka-band RF links", ACCENT_ORANGE),
]
for idx, (layer, desc, color) in enumerate(stack_layers):
    ly = Inches(1.6) + Inches(1.05) * idx
    card = add_card(slide, Inches(0.7), ly, Inches(5.6), Inches(0.9), border=color, fill=CARD_BG)
    add_textbox(slide, Inches(0.9), ly + Inches(0.05), Inches(5.2), Inches(0.35), layer, font_size=14, color=color, bold=True)
    add_textbox(slide, Inches(0.9), ly + Inches(0.42), Inches(5.2), Inches(0.4), desc, font_size=10, color=LIGHT_GRAY)
    if idx < len(stack_layers) - 1:
        add_textbox(slide, Inches(3.1), ly + Inches(0.85), Inches(1.0), Inches(0.25), "\u25bc", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

prio_headers = ["Priority", "Class", "Use Case"]
prio_rows = [
    ["P0", "EMERGENCY", "Life-critical alerts, abort commands"],
    ["P1", "URGENT", "Navigation corrections, hazard warnings"],
    ["P2", "NORMAL", "Telemetry, science data, commands"],
    ["P3", "LOW", "Software updates, bulk file transfer"],
    ["P4", "BULK", "Background sync, archival data"],
]
_d5 = [prio_headers] + prio_rows
add_table(slide, Inches(6.7), Inches(1.6), Inches(5.6), Inches(2.6), len(_d5), len(prio_headers), _d5, header_color=ACCENT_PURPLE)

snf_title_y = Inches(4.4)
add_textbox(slide, Inches(6.7), snf_title_y, Inches(5.6), Inches(0.35), "Store-and-Forward Steps", font_size=14, color=ACCENT_CYAN, bold=True)
add_accent_line(slide, Inches(6.7), snf_title_y + Inches(0.32), Inches(2.0), ACCENT_CYAN, Pt(2))
snf_steps = [
    "1. Bundle created at source with destination EID + TTL",
    "2. Node stores bundle in priority queue (buffer mgmt)",
    "3. Contact window opens \u2192 best next-hop selected by RL",
    "4. Bundle transmitted via convergence layer (LTP/TCPCL)",
    "5. Custody transfer confirms; repeat until destination",
]
for j, step in enumerate(snf_steps):
    add_textbox(slide, Inches(6.7), snf_title_y + Inches(0.4) + Inches(0.3) * j, Inches(5.6), Inches(0.3), step, font_size=10, color=LIGHT_GRAY)

standards_y = Inches(6.2)
standards_card = add_card(slide, Inches(0.7), standards_y, Inches(11.6), Inches(0.7), border=ACCENT_BLUE)
add_textbox(slide, Inches(0.9), standards_y + Inches(0.1), Inches(11.2), Inches(0.5), "Standards: RFC 9171 (BPv7)  \u2022  RFC 4838 (DTN Arch)  \u2022  RFC 5326 (LTP)  \u2022  CCSDS 734.2-B-1  \u2022  CCSDS 734.3-B-1 (SABR)", font_size=11, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9 — Network Topology
# ============================================================
print("Creating Slide 9: Network Topology...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "push")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "NETWORK TOPOLOGY"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "241 Nodes, 5 Tiers"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

tiers = [
    ("T1  Earth Ground", "6 nodes", "DSN stations: Goldstone, Madrid, Canberra", ACCENT_BLUE),
    ("T2  Earth Orbital", "51 nodes", "GEO relays + 48 LEO laser constellation", ACCENT_PURPLE),
    ("T3  Deep Space Transit", "4 nodes", "Lagrange point relays (ES-L4, ES-L5)", ACCENT_ORANGE),
    ("T4  Mars Orbital", "4 nodes", "Areostationary + polar orbit relays", ACCENT_RED),
    ("T5  Mars Surface", "176 nodes", "Bases, rovers, drones, sensor networks", YELLOW),
]

tier_y = Inches(1.6)
for tier_name, node_count, desc, color in tiers:
    bar = add_shape(
        slide, Inches(0.4), tier_y, Inches(0.08), Inches(0.65),
        fill_color=color, shape_type=MSO_SHAPE.RECTANGLE,
    )
    row_bg = add_shape(
        slide, Inches(0.48), tier_y, Inches(4.9), Inches(0.65),
        fill_color=CARD_BG, border_color=color, shape_type=MSO_SHAPE.RECTANGLE,
    )
    row_bg.line.color.rgb = color
    row_bg.line.width = Pt(1)

    tf = row_bg.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_top = Pt(2)
    p_name = tf.paragraphs[0]
    p_name.text = tier_name
    p_name.font.size = Pt(11)
    p_name.font.bold = True
    p_name.font.color.rgb = WHITE

    p_nodes = add_paragraph(tf)
    p_nodes.text = f"{node_count} \u2014 {desc}"
    p_nodes.font.size = Pt(8)
    p_nodes.font.color.rgb = LIGHT_GRAY

    tier_y += Inches(0.78)

add_image_safe(slide, os.path.join(DIAGRAMS_DIR, "5tier_network.png"), Inches(5.6), Inches(1.5), Inches(4.2), Inches(3.5))

callout_bg = add_shape(
    slide, Inches(0.4), Inches(5.7), Inches(9.2), Inches(0.65),
    fill_color=CARD_BG, border_color=ACCENT_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
)
callout_bg.line.color.rgb = ACCENT_BLUE
callout_bg.line.width = Pt(1)
ctf = callout_bg.text_frame
ctf.word_wrap = True
ctf.margin_left = Pt(12)
cp = ctf.paragraphs[0]
cp.text = "Multiple redundant paths \u00b7 No single point of failure \u00b7 Lagrange relays for conjunction coverage"
cp.font.size = Pt(11)
cp.font.color.rgb = ACCENT_CYAN
cp.alignment = PP_ALIGN.CENTER

add_footer(slide, 9, citations="[A2] AETHERIX topology.py (241 nodes, 5 tiers)  \u00b7  [1] NASA Deep Space Network (Goldstone/Madrid/Canberra)")


# ============================================================
# SLIDE 10 — 5-Tier Network Diagram
# ============================================================
print("Creating Slide 10: 5-Tier Network Diagram...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "wipe")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "5-TIER NETWORK DIAGRAM"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5))

add_image_safe(
    slide, os.path.join(DIAGRAMS_DIR, "5tier_network.png"),
    Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.3),
)

add_footer(slide, 10, citations="[A2] AETHERIX topology.py  \u00b7  [3] JPL Horizons (Lagrange ES-L4/L5 geometry)")


# ============================================================
# SLIDE 11 — Optical Communications
# ============================================================
print("Creating Slide 11: Optical Communications...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "OPTICAL COMMUNICATIONS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "10\u2013100x Faster Than RF"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

add_stat_card(slide, Inches(0.3), Inches(1.55), Inches(3.1), Inches(1.2), "100\u2013200", "Mbps at closest approach", GREEN)
add_stat_card(slide, Inches(3.55), Inches(1.55), Inches(3.1), Inches(1.2), "10\u201320", "Mbps average distance", ACCENT_BLUE)
add_stat_card(slide, Inches(6.8), Inches(1.55), Inches(3.1), Inches(1.2), "2\u20135", "Mbps at farthest point", ACCENT_RED)

eq_bullets = [
    "FSPL = 20\u00b7log\u2081\u2080(4\u03c0d/\u03bb)  \u2014 free-space path loss",
    "Gain = 10\u00b7log\u2081\u2080(\u03b7\u00b7(\u03c0D/\u03bb)\u00b2)  \u2014 antenna gain",
    "Pr = Pt + Gt + Gr \u2212 FSPL  \u2014 received power",
]
add_card(
    slide, Inches(0.3), Inches(2.9), Inches(4.8), Inches(1.8),
    "Key Equations",
    eq_bullets,
    border=ACCENT_BLUE,
)

config_headers = ["Parameter", "Value"]
config_rows = [
    ["Wavelength", "1550 nm"],
    ["Tx Power", "5 W"],
    ["Tx Aperture", "22 cm"],
    ["Rx Aperture", "1.0 m"],
]
_d3 = [config_headers] + config_rows
add_table(
    slide, Inches(5.3), Inches(2.9), Inches(4.4), Inches(1.8), len(_d3), len(config_headers), _d3,
    header_color=ACCENT_BLUE,
)

add_footer(slide, 11, citations="[5] CCSDS 141.0-B-1 (optical link)  \u00b7  [A4] AETHERIX link_budget.py  \u00b7  [1] NASA MRO (6 Mbps RF baseline)")


# ============================================================
# SLIDE 12 — Chart: Data Rate vs Distance
# ============================================================
print("Creating Slide 12: Data Rate vs Distance...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "data_rate_vs_distance.png"),
    "DATA RATE VS DISTANCE", "Optical Link Performance",
    "Data rate degrades with distance squared \u2014 200 Mbps at closest to 2 Mbps at farthest (10\u2013100\u00d7 capability over RF [1])", ACCENT_BLUE,
    citations="[5] CCSDS 141.0-B-1  \u00b7  [A4] AETHERIX link_budget.py (physics-derived capability)",
)


# ============================================================
# SLIDE 13 — Earth-Mars Journey
# ============================================================
print("Creating Slide 13: Earth-Mars Journey...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "push")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "EARTH-MARS JOURNEY"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "500 MB in 7 Hops"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

hop_headers = ["Hop", "Path", "Link"]
hop_rows = [
    ["1", "Rover \u2192 UHF Relay", "UHF S-band"],
    ["2", "UHF Relay \u2192 Areostationary", "UHF uplink"],
    ["3", "Areostationary \u2192 Polar Orbiter", "Crosslink"],
    ["4\u20135", "Polar Orbiter \u2192 LEO Constellation", "1550 nm laser"],
    ["6", "LEO Mesh \u2192 DSN Ground Station", "Optical downlink"],
    ["7", "DSN \u2192 Mission Operations Center", "TCP/IP fiber"],
]
_d4 = [hop_headers] + hop_rows
add_table(
    slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(2.7), len(_d4), len(hop_headers), _d4,
    header_color=ACCENT_BLUE,
)

add_stat_card(slide, Inches(0.2), Inches(4.4), Inches(2.35), Inches(1.1), "~13min", "End-to-end latency [3]", ACCENT_BLUE)
add_stat_card(slide, Inches(2.65), Inches(4.4), Inches(2.35), Inches(1.1), "<5%", "Protocol overhead", GREEN)
add_stat_card(slide, Inches(5.1), Inches(4.4), Inches(2.35), Inches(1.1), "98.7%", "Delivery (target)", ACCENT_PURPLE)
add_stat_card(slide, Inches(7.55), Inches(4.4), Inches(2.35), Inches(1.1), "7 hops", "Store-and-forward [A2]", ACCENT_ORANGE)

add_footer(slide, 13, citations="[A2] topology.py (7-hop path)  \u00b7  [3] JPL Horizons (light-time)  \u00b7  delivery ratio is a design target")


# ============================================================
# SLIDE 14 — RL Routing
# ============================================================
print("Creating Slide 14: RL Routing...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "wipe")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "RL-BASED ROUTING"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "AI Innovation"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

cgr_bullets = [
    "Pre-computed contact plans required",
    "Cannot adapt to unexpected link failures",
    "Computational cost grows O(n\u00b3) per route",
    "No learning from historical performance",
    "Single-path \u2014 ignores multipath opportunity",
]
add_card(
    slide, Inches(0.3), Inches(1.55), Inches(4.6), Inches(2.5),
    "Contact Graph Routing",
    cgr_bullets,
    border=ACCENT_RED,
)

rl_bullets = [
    "State: (node, neighbors, link quality, buffer)",
    "Actions: forward | store | drop | split",
    "Reward: R = \u03b1(delivery) \u2212 \u03b2(delay) \u2212 \u03b3(hops) \u2212 \u03b4(drops) \u2212 \u03b5(energy) [A1]",
    "\u03b5-greedy exploration, decay 0.995 (\u03b1=1.0, \u03b4=10.0) [A1]",
    "Q-table + federated aggregation across nodes",
]
add_card(
    slide, Inches(5.1), Inches(1.55), Inches(4.6), Inches(2.5),
    "RL Agent",
    rl_bullets,
    border=ACCENT_CYAN,
)

add_stat_card(slide, Inches(0.3), Inches(4.25), Inches(3.1), Inches(1.1), "713/800", "Forward decisions [A3]", GREEN)
add_stat_card(slide, Inches(3.55), Inches(4.25), Inches(3.1), Inches(1.1), "seconds", "Recovery vs hours [A1]", ACCENT_BLUE)
add_stat_card(slide, Inches(6.8), Inches(4.25), Inches(3.1), Inches(1.1), "Q-table", "Auditable policy [A1]", ACCENT_PURPLE)

add_footer(slide, 14, citations="[A1] rl_agent.py (reward fn, \u03b5-decay 0.995)  \u00b7  [A3] run_simulation Module 3 (training convergence 713/800)")


# ============================================================
# SLIDE 15 — Chart: RL Routing Heatmap
# ============================================================
print("Creating Slide 15: RL Routing Heatmap...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "rl_routing_heatmap.png"),
    "RL ROUTING HEATMAP", "Q-Value Distribution",
    "Learned routing preferences across network nodes \u2014 brighter = higher Q-value [A1]", ACCENT_CYAN,
    citations="[A1] rl_agent.py  \u00b7  [A3] run_simulation Module 3 (training)",
)


# ============================================================
# SLIDE 16 — Quantum Security
# ============================================================
print("Creating Slide 16: Quantum Security...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "cover")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "QUANTUM SECURITY"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Future-Proof"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

bb84_bullets = [
    "1. Alice sends random polarized photons (H/V, D/A bases) [13]",
    "2. Bob measures each photon in a random basis",
    "3. Public reconciliation \u2014 keep matching-basis bits",
    "4. Error estimation \u2014 sample subset for QBER",
    "5. Privacy amplification \u2014 universal hashing",
    "6. QBER < 11% \u2192 SECURE key established [15]",
]
add_card(
    slide, Inches(0.3), Inches(1.55), Inches(4.8), Inches(2.5),
    "BB84 Protocol",
    bb84_bullets,
    border=ACCENT_PURPLE,
)

pqc_bullets = [
    "Kyber (ML-KEM, FIPS 203) \u2014 key encapsulation [16]",
    "Dilithium (ML-DSA, FIPS 204) \u2014 digital signatures [17]",
    "NIST PQC standard \u2014 quantum-resistant",
    "Hybrid classical-quantum key exchange",
]
add_card(
    slide, Inches(5.3), Inches(1.55), Inches(4.4), Inches(2.5),
    "Post-Quantum Cryptography",
    pqc_bullets,
    border=ACCENT_ORANGE,
)

deploy_headers = ["Phase", "Protocol", "Key Rate"]
deploy_rows = [
    ["Phase 1 \u2014 LEO", "BB84", "1\u201310 kbps"],
    ["Phase 2 \u2014 GEO", "BB84 + E91", "10\u2013100 kbps"],
    ["Phase 3 \u2014 Mars", "E91 + Repeaters", "1+ kbps"],
]
_d5 = [deploy_headers] + deploy_rows
add_table(
    slide, Inches(0.3), Inches(4.25), Inches(9.4), Inches(1.5), len(_d5), len(deploy_headers), _d5,
    header_color=ACCENT_PURPLE,
)

add_footer(slide, 16, citations="[13] Bennett-Brassard 1984  \u00b7  [14] Ekert 1991  \u00b7  [15] Shor-Preskill 2000 (QBER<11%)  \u00b7  [16][17] NIST FIPS 203/204  \u00b7  [A5] qkd.py")


# ============================================================
# SLIDE 17 — Chart: QKD Security
# ============================================================
print("Creating Slide 17: QKD Security Chart...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "qkd_security.png"),
    "QKD SECURITY ANALYSIS", "Quantum Bit Error Rate",
    "QBER below 11% confirms secure key exchange \u2014 no eavesdropper detected [15]", ACCENT_PURPLE,
    citations="[15] Shor-Preskill 2000 (11% threshold)  \u00b7  [13] BB84  \u00b7  [A5] AETHERIX qkd.py",
)


# ============================================================
# SLIDE 18 — Orbital Mechanics
# ============================================================
print("Creating Slide 18: Orbital Mechanics...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "ORBITAL MECHANICS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Contact Windows"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

mars_headers = ["Parameter", "Value"]
mars_rows = [
    ["Semi-major axis", "1.524 AU"],
    ["Synodic period", "779.94 days"],
    ["Distance range", "54.6M \u2013 401M km"],
    ["Areostationary altitude", "17,032 km"],
    ["One-way light time", "3 \u2013 22 minutes"],
]
_d6 = [mars_headers] + mars_rows
add_table(
    slide, Inches(0.3), Inches(1.55), Inches(4.7), Inches(2.3), len(_d6), len(mars_headers), _d6,
    header_color=ACCENT_ORANGE,
)

window_headers = ["Window", "Availability", "Duration", "Data Rate"]
window_rows = [
    ["Optimal", "99%", "8\u201312 hrs", "100\u2013200 Mbps"],
    ["Good", "95%", "6\u20138 hrs", "20\u2013100 Mbps"],
    ["Fair", "85%", "4\u20136 hrs", "5\u201320 Mbps"],
    ["Blackout", "0%", "2\u20134 weeks", "\u2014"],
]
_d7 = [window_headers] + window_rows
add_table(
    slide, Inches(5.1), Inches(1.55), Inches(4.6), Inches(2.3), len(_d7), len(window_headers), _d7,
    header_color=ACCENT_BLUE,
)

add_footer(slide, 18, citations="[3] JPL Horizons (synodic 779.94 d, 54.6M\u2013401M km)  \u00b7  [A2] topology.py  \u00b7  [A4] link_budget.py (windows)")


# ============================================================
# SLIDE 19 — Data Flow Visual
# ============================================================
print("Creating Slide 19: Data Flow Visual...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "cover")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "DATA FLOW DIAGRAM"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

img_path_20 = os.path.join(DIAGRAMS_DIR, "data_flow.png")
add_image_safe(
    slide, img_path_20,
    Inches(0.4), Inches(1.2), Inches(9.2), Inches(5.5),
)

add_footer(slide, 19, citations="[A2] topology.py (full data path Mars\u2192Earth)")


# ============================================================
# SLIDE 20 — Performance
# ============================================================
print("Creating Slide 20: Performance...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "PERFORMANCE COMPARISON"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(9), Inches(0.35))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "AETHERIX vs Current"
p2.font.size, p2.font.color.rgb, p2.alignment = Pt(14), ACCENT_CYAN, PP_ALIGN.LEFT

perf_data = [
    ["Metric", "Current (MRO) [1]", "AETHERIX [A4]", "Note"],
    ["Downlink Rate", "0.5-6 Mbps", "2-200 Mbps", "10-100\u00d7 capability"],
    ["Daily Volume", "5-10 GB", "50-100 GB", "Est. (link-driven)"],
    ["Availability", "60-75%", ">95%", "Design target"],
    ["Routing", "Static (CGR)", "RL-adaptive", "Adaptive [A1]"],
    ["Security", "AES-256", "QKD + PQC", "Future-proof [16][17]"],
    ["Scalability", "5-10 assets", "241 nodes", "10-100\u00d7 [A2]"],
    ["Conjunction", "Blackout", "50-70% via L4/L5", "+50-70% [A8]"],
]
add_table(
    slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(3.4), len(perf_data), len(perf_data[0]) if perf_data else 0, perf_data,
    header_color=GREEN,
)

add_footer(slide, 20, citations="[1] NASA MRO (0.5\u20136 Mbps)  \u00b7  [A4] link_budget.py (2\u2013200 Mbps capability)  \u00b7  [A2] topology.py  \u00b7  [A8] Module 4  \u00b7  targets clearly labelled")


# ============================================================
# SLIDE 21 — Chart: Performance Comparison
# ============================================================
print("Creating Slide 21: Performance Comparison Chart...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "performance_comparison.png"),
    "PERFORMANCE COMPARISON", "AETHERIX vs Current Systems",
    "10\u2013100\u00d7 data-rate capability (link budget [A4] vs [1]) \u00b7 availability/routing are design targets", GREEN,
    citations="[1] NASA MRO  \u00b7  [A4] AETHERIX link_budget.py  \u00b7  [A2] topology.py",
)


# ============================================================
# SLIDE 22 — Implementation
# ============================================================
print("Creating Slide 22: Implementation...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "push")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "IMPLEMENTATION"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "What We Built"
p2.font.size, p2.font.color.rgb, p2.alignment = Pt(14), ACCENT_CYAN, PP_ALIGN.LEFT

impl_stats = [
    ("27", "Modules", GREEN),
    ("480", "Tests", ACCENT_BLUE),
    ("12", "Demos", ACCENT_PURPLE),
    ("5", "Policies", ACCENT_ORANGE),
]
stat_w = Inches(2.1)
stat_h = Inches(0.95)
stat_gap = Inches(0.25)
total_w_impl = 4 * stat_w + 3 * stat_gap
stat_x0_impl = (SLIDE_WIDTH - total_w_impl) // 2
for idx, (val, label, color) in enumerate(impl_stats):
    sx = stat_x0_impl + int((stat_w + stat_gap) * idx)
    add_stat_card(slide, sx, Inches(1.55), stat_w, stat_h, val, label, color)

modules_card = add_card(slide, Inches(0.5), Inches(2.8), Inches(4.3), Inches(1.7))
tf_mod = modules_card.text_frame
tf_mod.word_wrap = True
p_mod_title = tf_mod.paragraphs[0]
p_mod_title.text = "Core Modules"
p_mod_title.font.size, p_mod_title.font.bold, p_mod_title.font.color.rgb = Pt(13), True, ACCENT_CYAN
p_mod_title.space_after = Pt(6)
for mod_name in ["Link Budget", "RL Agent", "QKD", "Bundle Protocol", "Topology", "Simulation"]:
    p_mod = tf_mod.add_paragraph()
    p_mod.text = f"\u2022  {mod_name}"
    p_mod.font.size, p_mod.font.color.rgb = Pt(11), WHITE

std_card = add_card(slide, Inches(5.1), Inches(2.8), Inches(4.6), Inches(1.7))
tf_std = std_card.text_frame
tf_std.word_wrap = True
p_std_title = tf_std.paragraphs[0]
p_std_title.text = "Standards Compliance"
p_std_title.font.size, p_std_title.font.bold, p_std_title.font.color.rgb = Pt(13), True, ACCENT_CYAN
p_std_title.space_after = Pt(6)
standards_data = [
    ["Standard", "Title"],
    ["RFC 9171", "Bundle Protocol v7"],
    ["RFC 4838", "DTN Architecture"],
    ["RFC 5326", "LTP (deep space)"],
    ["RFC 9172", "BPSec (security)"],
    ["CCSDS 734.2-B-1", "CCSDS Bundle Protocol"],
    ["CCSDS 734.3-B-1", "SABR (contact routing)"],
    ["CCSDS 141.0-B-1", "Optical comms phys."],
]
std_tbl = add_table(
    slide, Inches(5.2), Inches(3.35), Inches(4.4), Inches(1.1), len(standards_data), len(standards_data[0]) if standards_data else 0, standards_data,
    header_color=ACCENT_BLUE,
)

add_footer(slide, 22, citations="[9][10][11][12] IETF RFCs  \u00b7  [2][5] CCSDS standards  \u00b7  [A1]\u2013[A8] AETHERIX modules (github.com/matx104/AETHERIX)")


# ============================================================
# SLIDE 23 — Roadmap
# ============================================================
print("Creating Slide 23: Roadmap...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "wipe")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "ROADMAP"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "From Demo to Deployment"
p2.font.size, p2.font.color.rgb, p2.alignment = Pt(14), ACCENT_CYAN, PP_ALIGN.LEFT

phases = [
    ("Phase 1-4  \u2713 Complete", "Topology \u2022 RL \u2022 QKD \u2022 Web UI \u2022 480 tests", GREEN),
    ("Phase 5  Network Sim", "ns-3 integration, realistic channel models", ACCENT_BLUE),
    ("Phase 6  Production", "DQN agent, ION-DTN stack, hardware-in-loop", ACCENT_PURPLE),
    ("Phase 7  Hardware", "SDR prototyping, optical ground station demo", ACCENT_ORANGE),
]
phase_w = Inches(4.6)
phase_h = Inches(0.7)
phase_gap = Inches(0.15)
phase_x = Inches(0.5)
for idx, (title, desc, color) in enumerate(phases):
    py = Inches(1.55) + int((phase_h + phase_gap) * idx)
    card = add_card(slide, phase_x, py, phase_w, phase_h)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    accent_bar = add_shape(
        slide, Inches(0.4), Inches(0.4), Inches(0.08), Inches(0.2), shape_type=MSO_SHAPE.RECTANGLE,
        fill_color=color,
    )
    tf_ph = card.text_frame
    tf_ph.word_wrap = True
    p_ph = tf_ph.paragraphs[0]
    p_ph.text = title
    p_ph.font.size, p_ph.font.bold, p_ph.font.color.rgb = Pt(12), True, color
    p_ph2 = tf_ph.add_paragraph()
    p_ph2.text = desc
    p_ph2.font.size, p_ph2.font.color.rgb = Pt(9), LIGHT_GRAY

add_footer(slide, 23, citations="[A1]\u2013[A8] AETHERIX source modules  \u00b7  [4] NASA DSOC (Psyche) heritage for optical  \u00b7  DQN/ns-3/ION-DTN on production roadmap")


# ============================================================
# SLIDE 24 — Conclusion
# ============================================================
print("Creating Slide 24: Conclusion...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "cover")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "CONCLUSION"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "AETHERIX Delivers"
p2.font.size, p2.font.color.rgb, p2.alignment = Pt(14), ACCENT_CYAN, PP_ALIGN.LEFT

concl_card_w = Inches(4.3)
concl_card_h = Inches(1.7)

problem_card = add_card(slide, Inches(0.5), Inches(1.6), concl_card_w, concl_card_h)
problem_card.line.color.rgb = ACCENT_RED
problem_card.line.width = Pt(2.5)
tf_prob = problem_card.text_frame
tf_prob.word_wrap = True
p_prob_t = tf_prob.paragraphs[0]
p_prob_t.text = "Problem Solved"
p_prob_t.font.size, p_prob_t.font.bold, p_prob_t.font.color.rgb = Pt(14), True, ACCENT_RED
p_prob_t.space_after = Pt(8)
p_prob_b = tf_prob.add_paragraph()
p_prob_b.text = "TCP/IP can\u2019t work at interplanetary distances"
p_prob_b.font.size, p_prob_b.font.color.rgb = Pt(12), WHITE

built_card = add_card(slide, Inches(5.2), Inches(1.6), concl_card_w, concl_card_h)
built_card.line.color.rgb = GREEN
built_card.line.width = Pt(2.5)
tf_built = built_card.text_frame
tf_built.word_wrap = True
p_built_t = tf_built.paragraphs[0]
p_built_t.text = "What We Built"
p_built_t.font.size, p_built_t.font.bold, p_built_t.font.color.rgb = Pt(14), True, GREEN
p_built_t.space_after = Pt(8)
p_built_b = tf_built.add_paragraph()
p_built_b.text = "BPv7 \u2022 RL Routing \u2022 QKD \u2022 Optical/RF Hybrid"
p_built_b.font.size, p_built_b.font.color.rgb = Pt(12), WHITE

conc_stats = [
    ("10-100\u00d7", "Data Rate [A4]", GREEN),
    (">95%", "Availability (target)", ACCENT_BLUE),
    ("RL", "Adaptive Routing [A1]", ACCENT_PURPLE),
    ("QKD", "Quantum Security [13][15]", ACCENT_ORANGE),
]
stat_w_c = Inches(2.1)
stat_h_c = Inches(1.1)
stat_gap_c = Inches(0.25)
total_w_c = 4 * stat_w_c + 3 * stat_gap_c
stat_x0_c = (SLIDE_WIDTH - total_w_c) // 2
for idx, (val, label, color) in enumerate(conc_stats):
    sx = stat_x0_c + int((stat_w_c + stat_gap_c) * idx)
    add_stat_card(slide, sx, Inches(3.7), stat_w_c, stat_h_c, val, label, color)

add_footer(slide, 24, citations="[1] NASA  \u00b7  [3] JPL Horizons  \u00b7  [9][12] IETF  \u00b7  [13][15] QKD  \u00b7  [A1][A2][A4][A5][A8] AETHERIX")


# ============================================================
# SLIDE — Radiation-Hardened Computing (LO-e)
# ============================================================
print("Creating slide: Radiation-Hardened Computing...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "push")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(12.0), Inches(0.7), "RADIATION-HARDENED COMPUTING", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_ORANGE, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(12.0), Inches(0.4), "Surviving the Space Radiation Environment (LO-e)", font_size=16, color=MED_GRAY)

eff_headers = ["Effect", "What it does", "Mitigation"]
eff_rows = [
    ["SEU (Single Event Upset)", "Bit flip in memory/register", "SECDED ECC"],
    ["MBU (Multiple Bit Upset)", "\u22652 adjacent bits flipped", "Bit interleaving"],
    ["SEL (Single Event Latchup)", "Parasitic short \u2014 destructive", "Current limit + power-cycle"],
    ["TID (Total Ionizing Dose)", "Cumulative degradation (krad)", "Rad-hard parts (RAD750) [18]"],
]
_d_rad = [eff_headers] + eff_rows
add_table(slide, Inches(0.3), Inches(1.6), Inches(7.0), Inches(2.4), len(_d_rad), len(eff_headers), _d_rad, header_color=ACCENT_ORANGE)

def_bullets = [
    "1. TMR \u2014 3 replicas + majority voter masks single fault",
    "2. SECDED (39,32) ECC \u2014 corrects 1, detects 2",
    "3. Memory scrubbing \u2014 rewrite before 2nd upset",
    "4. FDIR + watchdog \u2014 detect \u2192 isolate \u2192 reload",
    "5. SAFE-MODE after recovery budget exhausted",
]
add_card(slide, Inches(7.6), Inches(1.6), Inches(5.4), Inches(2.4), "Defense-in-Depth Stack", def_bullets, title_color=ACCENT_RED, border=ACCENT_RED, body_size=11)

add_textbox(slide, Inches(0.7), Inches(4.15), Inches(12.0), Inches(0.4), "DEMONSTRATED \u2014 run_simulation Module 6 [A6]", font_size=14, color=GREEN, bold=True)
rad_stats = [("37,159", "Raw upsets", ACCENT_ORANGE), ("200\u00d7", "Protection factor", GREEN), ("3,334\u00d7", "TMR gain (p=1e-4)", ACCENT_BLUE), ("2,127\u00d7", "RAD750 TID margin", ACCENT_PURPLE)]
for idx, (val, label, col) in enumerate(rad_stats):
    sx = Inches(0.3) + Inches(3.25) * idx
    add_stat_card(slide, sx, Inches(4.6), Inches(3.0), Inches(1.1), val, label, col, val_size=26)

add_textbox(slide, Inches(0.7), Inches(5.95), Inches(12.0), Inches(0.6), "Heritage: NASA RAD750 (Curiosity, Perseverance) [18]  \u00b7  ESA LEON3FT [19]  \u00b7  Code: src/computing/radiation.py [A6]  \u00b7  ~186 residual errors vs 37,159 raw (200\u00d7)", font_size=12, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_footer(slide, citations="[A6] AETHERIX radiation.py (demonstrated Module 6)  \u00b7  [18] BAE RAD750  \u00b7  [19] ESA LEON3FT")


# ============================================================
# SLIDE — Mission-Critical Data Prioritization (LO-f)
# ============================================================
print("Creating slide: Data Prioritization...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "cover")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(12.0), Inches(0.7), "MISSION-CRITICAL DATA PRIORITIZATION", font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(12.0), Inches(0.4), "Bandwidth Triage on a Starved, Intermittent Link (LO-f)", font_size=16, color=MED_GRAY)

pri_headers = ["Tier", "Class", "Examples", "BPv7"]
pri_rows = [
    ["P0", "Emergency / Safety", "Health, collision avoid, faults", "EMERGENCY"],
    ["P1", "Mission-critical", "Command ACKs, time-sensitive sci", "HIGH_SCIENCE"],
    ["P2", "High-priority", "Routine telemetry, scheduled sci", "STANDARD"],
    ["P4", "Low / Bulk", "Housekeeping, file xfer, SW images", "BULK"],
]
_d_pri = [pri_headers] + pri_rows
add_table(slide, Inches(0.3), Inches(1.6), Inches(7.4), Inches(2.4), len(_d_pri), len(pri_headers), _d_pri, header_color=ACCENT_CYAN)

lever_bullets = [
    "1. Compression (pre-transmit) [7][8]:",
    "   Telemetry \u2192 CCSDS 121.0-B-3 (Rice) \u2248 3\u00d7 [7]",
    "   Imagery \u2192 CCSDS 122.0-B-2 (wavelet) \u2248 10\u00d7 [8]",
    "2. Deadline-aware preemptive QoS [A7]:",
    "   strict priority \u2192 earliest deadline; defer if infeasible",
    "3. BPv7 fragmentation [9]:",
    "   send what fits; defer rest; link never idle",
]
add_card(slide, Inches(8.0), Inches(1.6), Inches(5.0), Inches(2.4), "Three Levers", lever_bullets, title_color=ACCENT_ORANGE, border=ACCENT_ORANGE, body_size=10)

add_textbox(slide, Inches(0.7), Inches(4.2), Inches(12.0), Inches(0.4), "DEMONSTRATED SCENARIO \u2014 src/routing/prioritization.py [A7]", font_size=14, color=GREEN, bold=True)
prio_stats = [("5/6", "Items delivered", GREEN), ("100%", "Link utilisation", ACCENT_CYAN), ("41%", "SW update fragmented", ACCENT_ORANGE), ("P0", "Emergency preempts", ACCENT_RED)]
for idx, (val, label, col) in enumerate(prio_stats):
    sx = Inches(0.3) + Inches(3.25) * idx
    add_stat_card(slide, sx, Inches(4.65), Inches(3.0), Inches(1.1), val, label, col, val_size=26)
add_footer(slide, citations="[A7] AETHERIX prioritization.py  \u00b7  [7] CCSDS 121.0-B-3  \u00b7  [8] CCSDS 122.0-B-2  \u00b7  [9] RFC 9171 (fragmentation)")


# ============================================================
# SLIDE — Trade-off Analysis
# ============================================================
print("Creating slide: Trade-off Analysis...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "wipe")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(12.0), Inches(0.7), "TRADE-OFF ANALYSIS \u2014 WHY THESE CHOICES", font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(12.0), Inches(0.4), "Every decision traded performance for auditability & reproducibility", font_size=15, color=MED_GRAY)

trade_headers = ["Decision", "Choice", "Rationale (vs alternative)"]
trade_rows = [
    ["Optical vs RF", "Hybrid 1550 nm + Ka-band fallback [4]", "10\u2013100\u00d7 throughput [A4]; RF survives clouds & corona"],
    ["Routing", "Custom Q-learning, not ION-DTN CGR [A1]", "Adapts to live state; CGR re-plans on stale schedule"],
    ["RL model", "Q-tables now, DQN later (Phase 6)", "Every Q-value human-auditable; trains in seconds"],
    ["State space", "Discretised, 241 nodes [A2]", "Right-sized for tabular policy; DQN path documented"],
    ["Reward weights", "\u03b1=1.0, \u03b4=10.0, \u03b5-decay 0.995 [A1]", "Drop penalty 10\u00d7 delivery to forbid bundle loss"],
]
_d_tr = [trade_headers] + trade_rows
add_table(slide, Inches(0.3), Inches(1.6), Inches(12.7), Inches(3.2), len(_d_tr), len(trade_headers), _d_tr, header_color=ACCENT_CYAN)

add_card(slide, Inches(0.3), Inches(5.0), Inches(12.7), Inches(1.4), "Bottom Line",
    ["Each choice sacrifices maximum theoretical performance for auditability and reproducibility \u2014 exactly what a defence",
     "and a research artefact require. DQN / ns-3 / ION-DTN are the documented production transition (Phases 7\u20139).",
     "DSOC heritage: NASA flew optical + RF side-by-side on Psyche [4] \u2014 AETHERIX mirrors that hybrid model."],
    title_color=GREEN, border=GREEN, body_size=12)
add_footer(slide, citations="[A1] rl_agent.py / training.py  \u00b7  [A4] link_budget.py  \u00b7  [4] NASA DSOC (Psyche)  \u00b7  [A2] topology.py")


# ============================================================
# SLIDE — Failure & Recovery
# ============================================================
print("Creating slide: Failure & Recovery...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(12.0), Inches(0.7), "FAILURE & RECOVERY", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_RED, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(12.0), Inches(0.4), "Autonomous Solar-Conjunction Link-Blackout Survival", font_size=15, color=MED_GRAY)
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12.0), Inches(0.35), "Scenario: Earth\u2013Sun\u2013Mars conjunction \u2014 corona collapses the 1550 nm link below the 0.3 forward threshold [A1]", font_size=11, color=LIGHT_GRAY)

path_headers = ["Path", "Band", "Quality", "Status", "Reward [A8]"]
path_rows = [
    ["Direct Mars \u2192 Earth", "1550 nm optical", "0.05", "CLOSED", "\u22121.438"],
    ["Mars \u2192 ES-L4 \u2192 Earth", "Ka-band RF", "0.65", "OPEN", "\u22120.201"],
    ["Mars \u2192 ES-L5 \u2192 Earth", "Ka-band RF", "0.60", "OPEN", "\u22120.201"],
]
_d_pa = [path_headers] + path_rows
add_table(slide, Inches(0.3), Inches(1.95), Inches(12.7), Inches(1.8), len(_d_pa), len(path_headers), _d_pa, header_color=ACCENT_RED)

rec_bullets = [
    "1. Detect \u2014 optical Q-value collapses (q<0.3)",
    "2. Re-route \u2014 agent picks highest-Q: ES-L4",
    "   (Ka-band RF, 60\u00b0 solar elongation, avoids corona)",
    "3. Prioritise \u2014 policy engine fires two rules:",
    "   P0 EMERGENCY \u2192 forward on best (Ka-band) link",
    "   P4 BULK \u2192 store locally, defer past conjunction",
]
add_card(slide, Inches(0.3), Inches(4.0), Inches(6.3), Inches(2.4), "How AETHERIX Recovers (automatically)", rec_bullets, title_color=ACCENT_CYAN, border=ACCENT_CYAN, body_size=11)

lag_bullets = [
    "ES-L4 / ES-L5 sit 60\u00b0 ahead/behind Earth.",
    "They keep line-of-sight to Mars around the solar",
    "limb even at true conjunction.",
    "Direct Earth\u2013Mars: 0% availability at conjunction.",
    "Via ES-L4/L5: 50\u201370% availability retained [A8].",
    "Geometry is Earth-side \u2014 no Mars relay solves this.",
]
add_card(slide, Inches(6.9), Inches(4.0), Inches(6.1), Inches(2.4), "Why Lagrange Relays [3][A2]", lag_bullets, title_color=ACCENT_PURPLE, border=ACCENT_PURPLE, body_size=11)
add_textbox(slide, Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.4), "Outcome: throughput drops (optical\u2192RF) but no mission-critical data lost. Run live: python run_simulation.py --module 4", font_size=11, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_footer(slide, citations="[A8] run_simulation Module 4 (\u22121.438 / \u22120.201)  \u00b7  [A1] rl_agent.py (0.3 threshold)  \u00b7  [3] JPL Horizons (Lagrange)  \u00b7  [A2] topology.py")


# ============================================================
# SLIDE — References (Industry & Scientific [1]-[20])
# ============================================================
print("Creating slide: References (Industry)...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(12.0), Inches(0.7), "REFERENCES", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(12.0), Inches(0.4), "Industry & Scientific Sources [1]\u2013[20]  \u00b7  three-layer attribution: [N] industry \u00b7 [AN] this project", font_size=13, color=MED_GRAY)

refs_left = ["[1] NASA JPL, Mars Relay Network User's Guide, 2024.",
             "[2] CCSDS, Bundle Protocol Spec., 734.2-B-1, 2021.",
             "[3] JPL Horizons, Earth\u2013Mars Ephemeris, 2025.",
             "[4] NASA, Deep Space Optical Comm. (DSOC), Psyche, 2024.",
             "[5] CCSDS, Optical Comm. Coding & Sync., 141.0-B-1, 2019.",
             "[6] CCSDS, TM Space Data Link Protocol, 131.0-B-3, 2017.",
             "[7] CCSDS, Lossless Data Compression, 121.0-B-3, 2020.",
             "[8] CCSDS, Image Data Compression, 122.0-B-2, 2017.",
             "[9] IETF, Bundle Protocol Version 7, RFC 9171, 2022.",
             "[10] IETF, Licklider Transmission Protocol, RFC 5326, 2008."]
refs_right = ["[11] IETF, DTN TCP Convergence Layer, RFC 7242, 2014.",
              "[12] IETF, Delay-Tolerant Networking Architecture, RFC 4838, 2007.",
              "[13] Bennett & Brassard, Quantum Cryptography, IEEE ICC, 1984.",
              "[14] Ekert, Quantum Cryptography Based on Bell's Theorem, PRL 67, 1991.",
              "[15] Shor & Preskill, Simple Proof of Security of BB84, PRL 85, 2000.",
              "[16] NIST, Module-Lattice-Based KEM, FIPS 203, 2024.",
              "[17] NIST, Module-Lattice-Based Digital Signature, FIPS 204, 2024.",
              "[18] BAE Systems, RAD750 Radiation-Hardened PowerPC, 2024.",
              "[19] ESA/Gaisler, LEON3FT Processor, 2023.",
              "[20] IETF, RFC 4838 \u00a73.1, 2007 (DTN challenges)."]
for j, ref in enumerate(refs_left):
    add_textbox(slide, Inches(0.4), Inches(1.6) + Inches(0.45) * j, Inches(6.3), Inches(0.4), ref, font_size=9, color=LIGHT_GRAY)
for j, ref in enumerate(refs_right):
    add_textbox(slide, Inches(6.9), Inches(1.6) + Inches(0.45) * j, Inches(6.3), Inches(0.4), ref, font_size=9, color=LIGHT_GRAY)
add_footer(slide, citations="Full URLs: github.com/matx104/AETHERIX  \u00b7  next slide: project sources [A1]\u2013[A8]")


# ============================================================
# SLIDE — References (Project [A1]-[A8])
# ============================================================
print("Creating slide: References (Project)...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(12.0), Inches(0.7), "REFERENCES (cont.)", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(12.0), Inches(0.4), "Project Sources [A1]\u2013[A8]  \u00b7  M. A. Tariq, AETHERIX, github.com/matx104/AETHERIX (2025)", font_size=13, color=MED_GRAY)

proj_left = ["[A1] RL Routing Agent \u2014 src/routing/rl_agent.py",
             "[A2] Network Topology \u2014 src/orbital/topology.py",
             "[A3] End-to-End Simulation \u2014 run_simulation.py Module 3 (RL training)",
             "[A4] Optical Link Budget Calculator \u2014 src/infrastructure/link_budget.py"]
proj_right = ["[A5] QKD Protocol Implementation \u2014 src/security/qkd.py",
              "[A6] Radiation Hardening Model \u2014 src/computing/radiation.py",
              "[A7] Data Prioritization Engine \u2014 src/routing/prioritization.py",
              "[A8] Failure & Recovery Simulation \u2014 run_simulation.py Module 4 (conjunction)"]
for j, ref in enumerate(proj_left):
    add_textbox(slide, Inches(0.4), Inches(1.65) + Inches(0.45) * j, Inches(6.3), Inches(0.4), ref, font_size=11, color=LIGHT_GRAY)
for j, ref in enumerate(proj_right):
    add_textbox(slide, Inches(6.9), Inches(1.65) + Inches(0.45) * j, Inches(6.3), Inches(0.4), ref, font_size=11, color=LIGHT_GRAY)

attr_bullets = ["Layer A \u2014 Industry/scientific baseline [1]\u2013[20] (e.g. NASA MRO 0.5\u20136 Mbps [1]).",
                "Layer B \u2014 This project's design decision [A1]\u2013[A8] (cite the code).",
                "Layer C \u2014 Demonstrated simulation results [A6][A8] (cite the run, e.g. radiation 200\u00d7 [A6]).",
                "Design targets (>95% availability, 2\u2013200 Mbps capability) are labelled as such, never as measured results."]
add_card(slide, Inches(0.4), Inches(3.7), Inches(12.8), Inches(2.2), "Three-Layer Attribution (resolves the prior citation gap)", attr_bullets, title_color=ACCENT_PURPLE, border=ACCENT_PURPLE, body_size=12)
add_footer(slide, citations="Reproducible: python run_simulation.py  \u00b7  480 tests: python -m pytest tests/ -q  \u00b7  matx104.github.io/AETHERIX")


# ============================================================
# SLIDE 25 — Thank You
# ============================================================
print("Creating Slide 25: Thank You...")
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

top_bar = add_shape(
    slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), shape_type=MSO_SHAPE.RECTANGLE,
    fill_color=ACCENT_BLUE,
)

ty_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.6), SLIDE_WIDTH - Inches(1), Inches(0.9)
)
p_ty = ty_box.text_frame.paragraphs[0]
p_ty.text = "THANK YOU"
p_ty.font.size, p_ty.font.bold, p_ty.font.color.rgb, p_ty.alignment = (
    Pt(64), True, WHITE, PP_ALIGN.CENTER,
)

add_accent_line(slide, Inches(3.6), Inches(1.7), Inches(2.8), ACCENT_CYAN)

q_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(2.0), SLIDE_WIDTH - Inches(1), Inches(0.6)
)
p_q = q_box.text_frame.paragraphs[0]
p_q.text = "Questions?"
p_q.font.size, p_q.font.bold, p_q.font.color.rgb, p_q.alignment = (
    Pt(36), True, ACCENT_CYAN, PP_ALIGN.CENTER,
)

ty_stats = [
    ("10-100\u00d7", "Data Rate [A4]", GREEN),
    (">95%", "Availability (target)", ACCENT_BLUE),
    ("RL", "Adaptive Routing [A1]", ACCENT_PURPLE),
    ("QKD", "Quantum Security [13][15]", ACCENT_ORANGE),
]
stat_w_ty = Inches(2.1)
stat_h_ty = Inches(0.95)
stat_gap_ty = Inches(0.25)
total_w_ty = 4 * stat_w_ty + 3 * stat_gap_ty
stat_x0_ty = (SLIDE_WIDTH - total_w_ty) // 2
for idx, (val, label, color) in enumerate(ty_stats):
    sx = stat_x0_ty + int((stat_w_ty + stat_gap_ty) * idx)
    add_stat_card(slide, sx, Inches(3.0), stat_w_ty, stat_h_ty, val, label, color)

contact_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(4.4), SLIDE_WIDTH - Inches(1), Inches(0.5)
)
p_contact = contact_box.text_frame.paragraphs[0]
p_contact.text = "Muhammad Abdullah Tariq"
p_contact.font.size, p_contact.font.bold, p_contact.font.color.rgb, p_contact.alignment = (
    Pt(20), True, WHITE, PP_ALIGN.CENTER,
)

links_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(4.9), SLIDE_WIDTH - Inches(1), Inches(0.4)
)
p_links = links_box.text_frame.paragraphs[0]
p_links.text = "matx104.github.io/AETHERIX  |  github.com/matx104/AETHERIX"
p_links.font.size, p_links.font.color.rgb, p_links.alignment = (
    Pt(12), MED_GRAY, PP_ALIGN.CENTER,
)

bottom_bar = add_shape(
    slide, Inches(0), SLIDE_HEIGHT - Inches(0.08), SLIDE_WIDTH, Inches(0.08), shape_type=MSO_SHAPE.RECTANGLE,
    fill_color=ACCENT_BLUE,
)


# ============================================================
# SAVE
# ============================================================

for _i, _sl in enumerate(prs.slides):
    if (_i + 1) in _SPEAKER_NOTES:
        _sl.notes_slide.notes_text_frame.text = _SPEAKER_NOTES[_i + 1]

pptx_path = os.path.join(OUTPUT_DIR, "AETHERIX_Presentation_Compact.pptx")
prs.save(pptx_path)
print(f"\n\u2713 Compact PPTX saved: {pptx_path}")
print(f"  Slides: {len(prs.slides)}")
