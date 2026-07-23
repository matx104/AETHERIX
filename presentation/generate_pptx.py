#!/usr/bin/env python3
"""
AETHERIX Animated Presentation Generator
Creates a professional PPTX with animations, charts, and diagrams.
"""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
import copy
from lxml import etree

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_DIR = os.path.join(BASE_DIR, "presentation", "output")
CHARTS_DIR = os.path.join(BASE_DIR, "visualizations", "charts")
DIAGRAMS_DIR = os.path.join(BASE_DIR, "visualizations", "diagrams")
IMG_DIR = os.path.join(BASE_DIR, "docs", "img")

os.makedirs(OUTPUT_DIR, exist_ok=True)

BG_DARK = RGBColor(0x0B, 0x0E, 0x1A)
BG_SLIDE = RGBColor(0x0D, 0x12, 0x22)
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xAA)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x4D, 0x4D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xCC)
MED_GRAY = RGBColor(0x6B, 0x7B, 0x96)
CARD_BG = RGBColor(0x14, 0x1B, 0x2D)
CARD_BORDER = RGBColor(0x1E, 0x2A, 0x42)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
YELLOW = RGBColor(0xFF, 0xD9, 0x3D)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text="", font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_card(slide, left, top, width, height, title="", body_lines=None, title_color=ACCENT_BLUE, body_color=LIGHT_GRAY, title_size=16, body_size=13, fill=CARD_BG, border=CARD_BORDER):
    shape = add_shape(slide, left, top, width, height, fill_color=fill, border_color=border)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Calibri"
    if body_lines:
        for line in body_lines:
            add_paragraph(tf, line, font_size=body_size, color=body_color, space_before=Pt(2), space_after=Pt(2))
    return shape


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    return add_shape(slide, left, top, width, height, fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    if os.path.exists(img_path):
        if width and height:
            return slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            return slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            return slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            return slide.shapes.add_picture(img_path, left, top)
    return None


def add_table(slide, left, top, width, height, rows, cols, data, header_color=ACCENT_BLUE, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = data[row_idx][col_idx] if row_idx < len(data) and col_idx < len(data[row_idx]) else ""
            cell.fill.solid()
            if row_idx == 0:
                cell.fill.fore_color.rgb = header_color
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                    p.font.size = Pt(12)
                    p.font.name = "Calibri"
                    p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 1 else RGBColor(0x18, 0x22, 0x36)
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = LIGHT_GRAY
                    p.font.size = Pt(11)
                    p.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    return table_shape


def add_entrance_animation(slide, shape, delay_ms=0, duration_ms=500, anim_type="fade"):
    slide_element = slide._element
    timing = slide_element.find(qn('p:timing'))
    if timing is None:
        timing = parse_xml(
            '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            '<p:childTnLst/>'
            '</p:cTn></p:par></p:tnLst>'
            '</p:timing>'
        )
        slide_element.append(timing)
    child_tn_lst = timing.find('.//' + qn('p:childTnLst'))
    shape_id = shape.shape_id
    seq_id = len(child_tn_lst) + 2
    if anim_type == "fade":
        anim_xml = f'<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><p:cTn id="{seq_id}" fill="hold"><p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{seq_id+1}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{seq_id+2}" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" nodeType="afterEffect"><p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id="{seq_id+3}" dur="{duration_ms}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set><p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{seq_id+4}" dur="{duration_ms}"/><p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl></p:cBhvr></p:animEffect></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
    else:
        anim_xml = f'<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="{seq_id}" fill="hold"><p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{seq_id+1}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{seq_id+2}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect"><p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id="{seq_id+3}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
    child_tn_lst.append(parse_xml(anim_xml))


def add_slide_transition(slide, transition_type="fade", duration_ms=700):
    slide_element = slide._element
    trans = {
        "fade": '<p:fade/>',
        "push": '<p:push dir="l"/>',
        "wipe": '<p:wipe dir="d"/>',
        "cover": '<p:cover dir="l"/>',
    }
    inner = trans.get(transition_type, '<p:fade/>')
    trans_xml = f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med" advClick="1">{inner}</p:transition>'
    existing = slide_element.find(qn('p:transition'))
    if existing is not None:
        slide_element.remove(existing)
    slide_element.insert(1, parse_xml(trans_xml))


# Total slide count shown in footers (title slide 1 + Thank-You have no footer).
TOTAL_SLIDES = 50
# Auto-incrementing footer counter so inserting slides never requires renumbering.
# Starts at 1 (the title slide is slide 1 and carries no footer); each footered
# slide is numbered in document order.
_footer_counter = [1]


def add_footer(slide, slide_num=None, total=None, citations=None):
    # slide_num is ignored (kept for call-site compatibility); numbering is
    # derived from document order so slides can be inserted freely.
    _footer_counter[0] += 1
    n = _footer_counter[0]
    if citations:
        add_textbox(slide, Inches(0.5), Inches(7.22), Inches(12.5), Inches(0.25),
                    citations, font_size=8, color=MED_GRAY)
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(5), Inches(0.4),
                "AETHERIX — Interplanetary Communication Network", font_size=10, color=MED_GRAY)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4),
                f"{n} / {TOTAL_SLIDES}", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


def add_stat_card(slide, x, y, w, h, value, label, value_color=ACCENT_BLUE, val_size=24):
    card = add_shape(slide, x, y, w, h, fill_color=CARD_BG, border_color=value_color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(val_size)
    p.font.color.rgb = value_color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(tf, label, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(2))
    return card


def new_slide():
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    return slide


def add_chart_slide(chart_file, title, subtitle, caption, accent_rgb, citations=None):
    slide = new_slide()
    add_slide_transition(slide, "fade")
    add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7),
                title, font_size=32, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.7), Inches(0.9), Inches(11.0), Inches(0.4),
                subtitle, font_size=16, color=accent_rgb)
    add_accent_line(slide, Inches(0.7), Inches(1.35), Inches(2.5), accent_rgb)
    add_image_safe(slide, chart_file, Inches(0.7), Inches(1.5), Inches(7.0), Inches(4.5))
    add_textbox(slide, Inches(0.7), Inches(6.2), Inches(11.0), Inches(0.3),
                caption, font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    add_footer(slide, citations=citations)
    return slide


_SPEAKER_NOTES = {
    1: "State your name clearly. Read the topic number and title exactly as on the exam paper. Pause to let examiners see it. Point to the logo. This is your first impression. (30 seconds)",
    2: "Quick overview of what we will cover. 13 topics across 29 slides. About 20 minutes. (20 seconds)",
    3: "This slide sets up the narrative arc. First explain what AETHERIX is in plain language - it's like the postal service for interplanetary space. Then pivot to the problem: TCP/IP was never designed for space. 22-minute delays break every assumption. Solar conjunction blackouts. Static routing. Vulnerable crypto. Each problem maps to one of our solutions. (1.5 minutes)",
    4: "Start with the scale. 54.6M to 401M km. Light itself takes 3-22 minutes one way. TCP/IP was designed for sub-second round trips. In space, by the time a packet acknowledgment returns, the link may be gone. Solar conjunction causes 2-week blackout. This is why NASA calls it Delay-Tolerant Networking. (1.5 minutes)",
    5: "Earth-Mars distance varies over the 780-day synodic period. (15 seconds)",
    6: "One-way light delay ranges from 3 to 22 minutes. (15 seconds)",
    7: "The key insight: instead of requiring an end-to-end connection like TCP, DTN works like the postal service. Each node takes custody of your data and forwards it when a link becomes available. Three pillars: BPv7 for the protocol, RL for intelligent routing, QKD for security. (1.5 minutes)",
    8: "Show the architecture. Six core modules feed into the simulation engine, which feeds the web showcase. Standards compliance at the bottom. Point to each module as you explain. (1 minute)",
    9: "Architecture diagram showing source modules feeding simulation engine and web demos.",
    10: "Distribution of 241 nodes across 5 network tiers. (15 seconds)",
    11: "BPv7 is the postal service protocol. Walk through the stack: Application at top, BPv7 as the store-and-forward layer, three convergence layers for different link types, physical at bottom. Custody transfer is the key innovation - each node takes legal responsibility. Priority P0 (emergency) to P4 (bulk). (2 minutes)",
    12: "Five priority classes from emergency to bulk. (15 seconds)",
    13: "Walk through the store-and-forward process. Bundle arrives, gets stored, node waits for next contact opportunity, then forwards. If link drops, bundle stays stored and retries. No data loss. This is fundamentally different from TCP's end-to-end retransmission. (1.5 minutes)",
    14: "241 nodes across 5 tiers. Walk through each tier. Earth Ground is the DSN - three stations around the globe for 24/7 coverage. Earth Orbital has LEO laser mesh for optical backhaul. Deep Space has Lagrange point relays - these are the critical innovation for conjunction coverage. Mars Orbital has areostationary relays at 17,032 km. Mars Surface is the most populated tier. (2 minutes)",
    15: "Visual overview of the 5-tier topology with 3 redundant paths.",
    16: "Network topology visualization.",
    17: "DSN ground station coverage map. (15 seconds)",
    18: "Key orbital relay positions in the network. (15 seconds)",
    19: "RUN THE LIVE DEMO from the Link Budget page. Show the 3 distance scenarios. 1550nm was chosen for telecom heritage and eye safety. FSPL at average distance is -365 dB. The telescope apertures are realistic for spacecraft. RF backup for reliability. (2 minutes)",
    20: "Data rate degrades with distance squared. (15 seconds)",
    21: "Link budget analysis showing free-space path loss as the dominant factor. (15 seconds)",
    22: "Interactive journey visualization.",
    23: "Protocol latency overhead comparison. (15 seconds)",
    24: "Daily data volume throughput comparison. (15 seconds)",
    25: "CGR is what NASA uses today. It's static - you have to pre-compute schedules. Our RL agent learns from experience. 8 state variables, 4 actions. The reward function balances delivery probability against delay, hops, drops, and energy. Multi-agent federated learning means agents at each node share knowledge. (2 minutes)",
    26: "RL agent Q-value heatmap showing learned routing preferences. (15 seconds)",
    27: "RL agent energy efficiency optimization. (15 seconds)",
    28: "BB84 is beautifully simple: send qubits, measure, compare bases, check QBER. If QBER is below 11%, no one listened in. CASCADE reconciliation and privacy amplification clean the key. E91 uses entanglement. Quantum repeaters at Lagrange points extend range. Post-quantum crypto as backup layer. (2 minutes)",
    29: "QKD security analysis with QBER threshold. (15 seconds)",
    30: "QKD key generation rate versus distance. (15 seconds)",
    31: "Mars and Earth dance around the Sun with a 26-month synodic period. Everything changes - distance, delay, bandwidth. At opposition we get great bandwidth. At conjunction, the Sun blocks everything. Our Lagrange relays at ES-L4 and ES-L5 maintain 50-70% capacity during conjunction. Doppler shift of 15 GHz at optical wavelengths requires real-time compensation. (1.5 minutes)",
    32: "Contact window prediction over the synodic period. (15 seconds)",
    33: "Space radiation is relentless. SEUs flip bits constantly - about 37,000 during a Mars transit. Our defense-in-depth: TMR masks logic faults (3,334x reliability gain), SECDED ECC corrects single-bit errors, scrubbing prevents double-bit accumulation, and FDIR with a watchdog catches everything else. The RAD750 can tolerate 200 krad - far above what a Mars mission needs. (1.5 minutes)",
    34: "Like an emergency room. P0 emergency gets sent immediately - it can even preempt an in-progress transfer. P1 mission-critical next. P2 routine science. P4 bulk data fills remaining bandwidth. Compression multiplies effective capacity: 3x for telemetry, 10x for images, 50x for video. Our scheduler keeps the link at 100% utilization by fragmenting large bundles. (1.5 minutes)",
    35: "Walk through the 7-hop journey. 500MB from Perseverance to JPL. Total transit ~13 min vs 12.5 min light-time - near speed of light! DTN overhead under 5%. Key point: if link drops at hop 5, the bundle stays stored at hop 4 and retries. No data loss. RUN LIVE DEMO if time permits. (2 minutes)",
    36: "End-to-end bundle journey through all protocol layers.",
    37: "Visual data flow through the protocol stack.",
    38: "Be honest. Left column is simulated: 241-node topology, RL routing converges in 140 episodes, autonomous conjunction reroute, BB84 detection, 200x radiation reduction. Right column is design targets: 2-200 Mbps from link model, >95% availability and cost are goals not yet measured. (1 minute)",
    39: "Side-by-side performance comparison chart. (15 seconds)",
    40: "Optical versus RF link capability radar chart. (15 seconds)",
    41: "This is real, working code. 27 Python modules, 480 tests, 12 interactive demos. All the physics is real - no mocked data. The showcase site has live calculators you can use right now. Standards compliance is complete - CCSDS, IETF, and NIST. (1.5 minutes)",
    42: "Bandwidth evolution over deep-space missions. (15 seconds)",
    43: "Mission timeline deployment phases. (15 seconds)",
    44: "Phases 1-4 are done - this is what you see today. Phase 5: ns-3 simulation for realistic network modeling. Phase 6: Upgrade to DQN and integrate with NASA's ION-DTN implementation. Phase 7: Hardware prototypes with SDRs and optical links. (1.5 minutes)",
    45: "Summarize the problem and solution clearly. Re-read the exam topic verbatim. Point to the numbers. Offer to show live demos or answer questions. Thank the examiners. (1 minute)",
    46: "Summarize: simulated achievements are 241-node topology, RL routing, QKD security, 200x radiation reduction. Design targets are 10-100x data rate and >95% availability. Invite questions confidently. Thank the audience. (30 seconds)",
}

# ============================================================
# SLIDES 1-46
# ============================================================

# --- SLIDE 1 — Introduction (Title Hero) ---
print("Creating Slide 1: Introduction...")
slide = new_slide()
add_slide_transition(slide, "fade")
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.RECTANGLE)
add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.5), "AETHERIX", font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(4.5), Inches(2.5), Inches(4.3), ACCENT_CYAN, Pt(4))
add_textbox(slide, Inches(1.0), Inches(2.8), Inches(11.3), Inches(1.0), "Autonomous Extraterrestrial High-throughput Enhancing Routing\nand Inter-planetary eXchange", font_size=22, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.0), Inches(4.0), Inches(9.3), Inches(0.5), "EduQual Level 6 Diploma in AI Operations  |  Topic 59", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.0), Inches(4.5), Inches(9.3), Inches(0.5), "Building an Interplanetary Communication Network with DTN,\nQuantum Communication, and Space-Based Infrastructure for Mars Mission Support", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(3.5), Inches(5.3), Inches(6.3), ACCENT_BLUE, Pt(2))
add_textbox(slide, Inches(2.0), Inches(5.6), Inches(9.3), Inches(0.5), "Muhammad Abdullah Tariq", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.0), Inches(6.1), Inches(9.3), Inches(0.5), "September 2026", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.0), Inches(6.7), Inches(9.3), Inches(0.4), "matx104.github.io/AETHERIX  |  github.com/matx104/AETHERIX", font_size=12, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

# --- SLIDE 2 — Agenda ---
print("Creating Slide 2: Agenda...")
slide = new_slide()
add_slide_transition(slide, "push")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "PRESENTATION AGENDA", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_footer(slide, 2, citations="[A2] topology.py (241 nodes)  ·  References slides at end of deck")

agenda_items = [
    ("01", "The Challenge", "Why space breaks the internet", ACCENT_BLUE),
    ("02", "Architecture", "DTN + AI + Quantum Security", ACCENT_CYAN),
    ("03", "DTN & BPv7", "Store-and-forward foundation", ACCENT_PURPLE),
    ("04", "Topology", "241 nodes across two worlds", ACCENT_BLUE),
    ("05", "Link Budget", "1550nm laser analysis", ACCENT_ORANGE),
    ("06", "RL Routing", "Multi-agent federated Q-learning", ACCENT_CYAN),
    ("07", "Quantum Security", "BB84/E91 + repeater chains", ACCENT_PURPLE),
    ("08", "Orbital Mechanics", "Contact windows & synodic period", ACCENT_BLUE),
    ("09", "Radiation Hardening", "SEU/TID, TMR, ECC, scrubbing, FDIR", ACCENT_RED),
    ("10", "Data Prioritization", "QoS triage, compression, preemption", ACCENT_ORANGE),
    ("11", "Mars Mission", "End-to-end simulation walkthrough", ACCENT_ORANGE),
    ("12", "Performance", "AETHERIX vs current systems", GREEN),
    ("13", "Roadmap", "CCSDS, IETF, deployment phases", ACCENT_CYAN),
    ("14", "Q&A", "Summary and live demo", WHITE),
]

_per_col = 7
for i, (num, title, desc, color) in enumerate(agenda_items):
    col_x = Inches(0.7) if i < _per_col else Inches(6.8)
    row_y = Inches(1.25) + Inches(0.80) * (i % _per_col)
    card = add_card(slide, col_x, row_y, Inches(5.8), Inches(0.72), border=color)
    add_textbox(slide, col_x + Inches(0.15), row_y + Inches(0.06), Inches(5.5), Inches(0.32), f"{num}  {title}", font_size=13, color=WHITE, bold=True)
    add_textbox(slide, col_x + Inches(0.15), row_y + Inches(0.38), Inches(5.5), Inches(0.3), desc, font_size=10, color=LIGHT_GRAY)
    add_entrance_animation(slide, card, delay_ms=80 * i, anim_type="fade")

# --- SLIDE 3 — What is AETHERIX ---
print("Creating Slide 3: What is AETHERIX...")
slide = new_slide()
add_slide_transition(slide, "fade")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "WHAT IS AETHERIX?", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11.0), Inches(0.4), "Overview & The Problem", font_size=16, color=MED_GRAY)
add_footer(slide, 3, citations="[1] NASA MRN 2024  ·  [3] JPL Horizons  ·  [12] RFC 4838  ·  [A2] topology.py")

card_left = add_card(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(3.8), border=ACCENT_BLUE)
add_textbox(slide, Inches(0.9), Inches(1.7), Inches(5.2), Inches(0.45), "What is AETHERIX?", font_size=18, color=ACCENT_BLUE, bold=True)
add_accent_line(slide, Inches(0.9), Inches(2.1), Inches(2.0), ACCENT_BLUE, Pt(2))
what_lines = [
    "Autonomous Extraterrestrial High-throughput Enhancing\nRouting and Inter-planetary eXchange",
    "An architecture for delay-tolerant networking (DTN)\nbetween Earth and Mars",
    "Replaces static Contact Graph Routing with\nreinforcement-learning-based adaptive routing",
    "Secures links via Quantum Key Distribution\n(BB84 & E91 protocols with repeater chains)",
    "Hybrid optical/RF communications across\na 5-tier, 241-node interplanetary topology",
    "Open-source, standards-compliant (CCSDS, IETF)\nresearch platform for deep-space networking",
]
for j, line in enumerate(what_lines):
    add_textbox(slide, Inches(0.9), Inches(2.2) + Inches(0.5) * j, Inches(5.2), Inches(0.5), line, font_size=10, color=LIGHT_GRAY)

card_right = add_card(slide, Inches(6.7), Inches(1.6), Inches(5.6), Inches(3.8), border=ACCENT_RED)
add_textbox(slide, Inches(6.9), Inches(1.7), Inches(5.2), Inches(0.45), "The Problem", font_size=18, color=ACCENT_RED, bold=True)
add_accent_line(slide, Inches(6.9), Inches(2.1), Inches(2.0), ACCENT_RED, Pt(2))
problem_lines = [
    "Earth–Mars distance varies from 54.6M to 401M km\n— a 7x range swing every 780 days",
    "One-way light delay is 3–22 minutes;\nround-trip TCP ACK is 6–44 minutes",
    "Solar conjunction causes ~2-week communication\nblackouts twice per synodic period",
    "Deep-space links are inherently intermittent:\nno continuous end-to-end path exists",
    "Static routing cannot adapt to dynamic\ncontact schedules and link degradation",
    "Classical key exchange is infeasible across\ninterplanetary distances",
]
for j, line in enumerate(problem_lines):
    add_textbox(slide, Inches(6.9), Inches(2.2) + Inches(0.5) * j, Inches(5.2), Inches(0.5), line, font_size=10, color=LIGHT_GRAY)

callout_shape = add_shape(slide, Inches(0.7), Inches(5.7), Inches(11.6), Inches(0.9), fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(1.0), Inches(5.8), Inches(11.0), Inches(0.7), "AETHERIX addresses every one of these challenges — harnessing DTN store-and-forward, AI-driven adaptive routing, quantum-secured key distribution, and hybrid optical/RF links to enable reliable interplanetary communication.", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# --- SLIDE 4 — The Distance ---
print("Creating Slide 4: The Distance...")
slide = new_slide()
add_slide_transition(slide, "push")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "THE DISTANCE", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_RED, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11.0), Inches(0.4), "Why Space Breaks the Internet", font_size=16, color=MED_GRAY)
add_footer(slide, 4, citations="[3] JPL Horizons (distance/light-time)  ·  [12] RFC 4838  ·  [1] NASA MRO data rate")

left_headers = ["Parameter", "Value"]
left_rows = [
    ["Closest approach", "54.6 million km"],
    ["Farthest distance", "401 million km"],
    ["Distance variation", "7.3× (735%)"],
    ["One-way light time", "3 – 22 minutes"],
    ["Synodic period", "780 days (26 months)"],
    ["Solar conjunction blackout", "~2 weeks (twice/period)"],
]
_d1 = [left_headers] + left_rows
add_table(slide, Inches(0.7), Inches(1.5), Inches(5.6), Inches(3.2), len(_d1), len(left_headers), _d1, header_color=ACCENT_RED)

right_headers = ["TCP/IP Assumption", "Deep-Space Reality"]
right_rows = [
    ["Low latency (< 200 ms)", "6 – 44 min round trip"],
    ["Continuous connectivity", "Scheduled contact windows only"],
    ["End-to-end path exists", "No simultaneous path possible"],
    ["Reliable ACKs in seconds", "ACKs take minutes to hours"],
]
_d2 = [right_headers] + right_rows
add_table(slide, Inches(6.7), Inches(1.5), Inches(5.6), Inches(2.2), len(_d2), len(right_headers), _d2, header_color=ACCENT_ORANGE)

quote_card = add_card(slide, Inches(6.7), Inches(4.0), Inches(5.6), Inches(1.3), border=ACCENT_ORANGE)
add_textbox(slide, Inches(6.9), Inches(4.1), Inches(5.2), Inches(1.1), '"TCP/IP assumes the network is fast, reliable, and always connected.\nDeep space is none of those things."', font_size=12, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

stat_y = Inches(5.6)
stat_labels = ["54.6M km", "401M km", "780 days", "0.5–6 Mbps"]
stat_sublabels = ["Min Distance", "Max Distance", "Synodic Period", "Current DSN Rate"]
stat_colors = [ACCENT_BLUE, ACCENT_RED, ACCENT_ORANGE, ACCENT_CYAN]
for k in range(4):
    sx = Inches(0.7) + Inches(3.0) * k
    add_shape(slide, sx, stat_y, Inches(2.7), Inches(1.1), fill_color=CARD_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_accent_line(slide, sx, stat_y, Inches(2.7), stat_colors[k], Pt(3))
    add_textbox(slide, sx + Inches(0.15), stat_y + Inches(0.15), Inches(2.4), Inches(0.45), stat_labels[k], font_size=22, color=stat_colors[k], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, sx + Inches(0.15), stat_y + Inches(0.6), Inches(2.4), Inches(0.35), stat_sublabels[k], font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# --- SLIDE 5 — Chart: Distance Over Time ---
print("Creating Slide 5: Chart — Distance Over Time...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "distance_over_time.png"),
    "DISTANCE OVER TIME", "Earth-Mars Distance Variation",
    "Earth-Mars distance over the 780-day synodic period", ACCENT_RED, citations="[3] JPL Horizons (synodic 779.94 d)  ·  [1] NASA"
)

# --- SLIDE 6 — Chart: Light Time Delay ---
print("Creating Slide 6: Chart — Light Time Delay...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "light_time_delay.png"),
    "LIGHT TIME DELAY", "One-Way Signal Propagation Time",
    "Signal delay ranges from 3 to 22 minutes depending on orbital position", ACCENT_ORANGE, citations="[3] JPL Horizons (c=299,792 km/s)  ·  [12] RFC 4838"
)

# --- SLIDE 7 — The Answer ---
print("Creating Slide 7: The Answer...")
slide = new_slide()
add_slide_transition(slide, "wipe")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "THE ANSWER", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11.0), Inches(0.4), "Delay-Tolerant Networking", font_size=16, color=MED_GRAY)
add_footer(slide, 5, citations="[9] RFC 9171 BPv7  ·  [10] RFC 5326 LTP  ·  [11] RFC 7242 TCPCL  ·  [12] RFC 4838")

flow_labels = ["Bundle", "Store", "Wait", "Forward", "Deliver"]
flow_colors = [ACCENT_BLUE, ACCENT_PURPLE, ACCENT_ORANGE, ACCENT_CYAN, GREEN]
flow_descs = [
    "Create data\nbundle (BPv7)",
    "Buffer at\nlocal node",
    "Hold until\ncontact window",
    "Transmit to\nnext hop",
    "Reach final\ndestination",
]
card_w = Inches(2.0)
card_h = Inches(1.6)
start_x = Inches(0.5)
gap = Inches(0.35)
arrow_w = Inches(0.3)
flow_y = Inches(1.6)

for idx, (label, color, desc) in enumerate(zip(flow_labels, flow_colors, flow_descs)):
    cx = start_x + (card_w + arrow_w + gap) * idx
    card = add_card(slide, cx, flow_y, card_w, card_h, border=color)
    add_textbox(slide, cx + Inches(0.1), flow_y + Inches(0.1), card_w - Inches(0.2), Inches(0.35), label, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, cx + Inches(0.1), flow_y + Inches(0.5), card_w - Inches(0.2), Inches(0.9), desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_entrance_animation(slide, card, delay_ms=150 * idx, anim_type="fade")
    if idx < len(flow_labels) - 1:
        arrow_x = cx + card_w + Inches(0.05)
        add_textbox(slide, arrow_x, flow_y + Inches(0.5), arrow_w, Inches(0.5), "\u2192", font_size=28, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

stat_card_data = [
    ("BPv7", "Bundle Protocol v7\nRFC 9171 \u2014 store &\nforward with custody", ACCENT_BLUE),
    ("AI Routing", "Multi-agent federated\nQ-learning adapts in\nreal time to contacts", ACCENT_PURPLE),
    ("QKD", "BB84 & E91 quantum\nkey distribution with\nrepeater chains", ACCENT_CYAN),
]
stat_y = Inches(3.5)
for idx, (title, desc, color) in enumerate(stat_card_data):
    sx = Inches(0.7) + Inches(4.0) * idx
    add_stat_card(slide, sx, stat_y, Inches(3.6), Inches(1.4), title, desc, value_color=color)

postal_card = add_card(slide, Inches(0.7), Inches(5.2), Inches(11.6), Inches(1.0), border=ACCENT_BLUE)
add_textbox(slide, Inches(1.0), Inches(5.3), Inches(11.0), Inches(0.8), 'Think of it like the postal service: you don\'t need a continuous connection between sender and receiver \u2014 each post office (node) stores the letter (bundle) until the next truck (contact window) is available, then forwards it one hop closer to the destination.', font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# --- SLIDE 8 — System Architecture ---
print("Creating Slide 8: System Architecture...")
slide = new_slide()
add_slide_transition(slide, "cover")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "SYSTEM ARCHITECTURE", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11.0), Inches(0.4), "Five Core Modules", font_size=16, color=MED_GRAY)
add_footer(slide, 6, citations="[A1] rl_agent.py  ·  [A2] topology.py  ·  [A3] simulator.py  ·  [A4] link_budget.py  ·  [A5] qkd.py")

modules = [
    ("1", "Infrastructure", "Link budget & RF/optical analysis", ACCENT_BLUE),
    ("2", "Routing", "RL agent, BPv7, store-and-forward", ACCENT_PURPLE),
    ("3", "Security", "QKD (BB84/E91) & repeater chains", ACCENT_CYAN),
    ("4", "Orbital", "Contact windows & Doppler analysis", ACCENT_ORANGE),
    ("5", "Simulation", "Full network simulation engine", GREEN),
]

for idx, (num, name, desc, color) in enumerate(modules):
    my = Inches(1.6) + Inches(0.95) * idx
    circle = add_shape(slide, Inches(0.7), my, Inches(0.6), Inches(0.6), fill_color=color, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide, Inches(0.7), my + Inches(0.08), Inches(0.6), Inches(0.45), num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), my + Inches(0.0), Inches(4.0), Inches(0.35), name, font_size=16, color=color, bold=True)
    add_textbox(slide, Inches(1.5), my + Inches(0.35), Inches(4.0), Inches(0.35), desc, font_size=11, color=LIGHT_GRAY)

table_headers = ["Module", "Engine", "Showcase"]
table_rows = [
    ["Infrastructure", "OpticalLinkBudget + RFLinkBudget", "Earth-Mars link analysis"],
    ["Routing", "RLRoutingAgent + ForwardingEngine", "Multi-hop DTN delivery"],
    ["Security", "BB84Protocol + E91Protocol", "Quantum-secured key exchange"],
    ["Orbital", "ContactWindows + Doppler", "Synodic contact prediction"],
    ["Simulation", "Simulator + PolicyEngine", "End-to-end scenario runs"],
]
_d4 = [table_headers] + table_rows
add_table(slide, Inches(6.4), Inches(1.6), Inches(6.0), Inches(4.5), len(_d4), len(table_headers), _d4, header_color=ACCENT_BLUE)

# --- SLIDE 9 — Architecture Diagram ---
print("Creating Slide 9: Architecture Diagram...")
slide = new_slide()
add_slide_transition(slide, "fade")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "ARCHITECTURE DIAGRAM", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_footer(slide, 7, citations="[A2] AETHERIX topology.py  ·  github.com/matx104/AETHERIX")

arch_img_path = os.path.join(DIAGRAMS_DIR, "system_architecture.png")
add_image_safe(slide, arch_img_path, Inches(0.7), Inches(1.2), Inches(11.6), Inches(5.8))

# --- SLIDE 10 — Chart: Network Tier Distribution ---
print("Creating Slide 10: Chart \u2014 Network Tier Distribution...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "network_tier_distribution.png"),
    "NETWORK TIER DISTRIBUTION", "241 Nodes Across 5 Tiers",
    "Node count by network tier \u2014 Mars Surface is the most populated", ACCENT_BLUE, citations="[A2] AETHERIX topology.py (241 nodes, 5 tiers)"
)

# --- SLIDE 11 — BPv7 Deep Dive ---
print("Creating Slide 11: BPv7 Deep Dive...")
slide = new_slide()
add_slide_transition(slide, "push")
add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11.0), Inches(0.7), "BPv7 DEEP DIVE", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.7), Inches(0.95), Inches(3.0), ACCENT_CYAN, Pt(3))
add_textbox(slide, Inches(0.7), Inches(1.05), Inches(11.0), Inches(0.4), "The Foundation", font_size=16, color=MED_GRAY)
add_footer(slide, 8, citations="[9] RFC 9171  ·  [2] CCSDS 734.2-B-1  ·  [10] RFC 5326 LTP  ·  [11] RFC 7242 TCPCL  ·  [12] RFC 4838")

stack_layers = [
    ("Application Layer", "User data, commands, telemetry", ACCENT_BLUE),
    ("Bundle Protocol v7", "Routing, custody, fragmentation, TTL", ACCENT_PURPLE),
    ("Convergence Layer", "LTP (space), TCPCL (Earth), UDPCL (ISL)", ACCENT_CYAN),
    ("Physical Layer", "Optical 1550nm / Ka-band RF links", ACCENT_ORANGE),
]
for idx, (layer, desc, color) in enumerate(stack_layers):
    ly = Inches(1.6) + Inches(1.05) * idx
    card = add_card(slide, Inches(0.7), ly, Inches(5.6), Inches(0.9), border=color, fill=CARD_BG)
    add_textbox(slide, Inches(0.9), ly + Inches(0.05), Inches(5.2), Inches(0.35), layer, font_size=14, color=color, bold=True)
    add_textbox(slide, Inches(0.9), ly + Inches(0.42), Inches(5.2), Inches(0.4), desc, font_size=10, color=LIGHT_GRAY)
    if idx < len(stack_layers) - 1:
        add_textbox(slide, Inches(3.1), ly + Inches(0.85), Inches(1.0), Inches(0.25), "\u25bc", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

prio_headers = ["Priority", "Class", "Use Case"]
prio_rows = [
    ["P0", "EMERGENCY", "Life-critical alerts, abort commands"],
    ["P1", "URGENT", "Navigation corrections, hazard warnings"],
    ["P2", "NORMAL", "Telemetry, science data, commands"],
    ["P3", "LOW", "Software updates, bulk file transfer"],
    ["P4", "BULK", "Background sync, archival data"],
]
_d5 = [prio_headers] + prio_rows
add_table(slide, Inches(6.7), Inches(1.6), Inches(5.6), Inches(2.6), len(_d5), len(prio_headers), _d5, header_color=ACCENT_PURPLE)

snf_title_y = Inches(4.4)
add_textbox(slide, Inches(6.7), snf_title_y, Inches(5.6), Inches(0.35), "Store-and-Forward Steps", font_size=14, color=ACCENT_CYAN, bold=True)
add_accent_line(slide, Inches(6.7), snf_title_y + Inches(0.32), Inches(2.0), ACCENT_CYAN, Pt(2))
snf_steps = [
    "1. Bundle created at source with destination EID + TTL",
    "2. Node stores bundle in priority queue (buffer mgmt)",
    "3. Contact window opens \u2192 best next-hop selected by RL",
    "4. Bundle transmitted via convergence layer (LTP/TCPCL)",
    "5. Custody transfer confirms; repeat until destination",
]
for j, step in enumerate(snf_steps):
    add_textbox(slide, Inches(6.7), snf_title_y + Inches(0.4) + Inches(0.3) * j, Inches(5.6), Inches(0.3), step, font_size=10, color=LIGHT_GRAY)

standards_y = Inches(6.2)
standards_card = add_card(slide, Inches(0.7), standards_y, Inches(11.6), Inches(0.7), border=ACCENT_BLUE)
add_textbox(slide, Inches(0.9), standards_y + Inches(0.1), Inches(11.2), Inches(0.5), "Standards: RFC 9171 (BPv7)  \u2022  RFC 4838 (DTN Arch)  \u2022  RFC 5326 (LTP)  \u2022  CCSDS 734.2-B-1 (Bundle Protocol)  \u2022  CCSDS 734.3-B-1 (SABR)", font_size=11, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# --- SLIDE 12 — Chart: Bundle Priority Classes ---
print("Creating Slide 12: Chart \u2014 Bundle Priority Classes...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "bundle_priority_classes.png"),
    "BUNDLE PRIORITY CLASSES", "BPv7 QoS Tiers",
    "Five priority levels from P0 Emergency to P4 Bulk ensure critical data arrives first", ACCENT_PURPLE, citations="[9] RFC 9171  ·  [2] CCSDS 734.2-B-1"
)

# ── SLIDE 13 ── DTN Store-and-Forward ──────────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "DTN STORE-AND-FORWARD"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "How It Works"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

tcp_bullets = [
    "Requires continuous end-to-end connection",
    "Round-trip time > 40 minutes breaks ACKs",
    "Packet loss triggers aggressive retransmission",
    "Congestion control collapses at interplanetary scale",
]
add_card(
    slide, Inches(0.4), Inches(1.6), Inches(4.4), Inches(2.8),
    "TCP/IP in Space \u2014 Fails",
    tcp_bullets,
    border=ACCENT_RED,
)

dtn_bullets = [
    "Store-and-forward: custody transfer at each hop",
    "Bundle Protocol v7 (RFC 9171) \u2014 no end-to-end session",
    "Opportunistic contacts exploited automatically",
    "Tolerates hours-to-days of link disruption",
]
add_card(
    slide, Inches(5.2), Inches(1.6), Inches(4.4), Inches(2.8),
    "DTN Works",
    dtn_bullets,
    border=ACCENT_CYAN,
)

ltp_bullets = [
    "Licklider Transmission Protocol",
    "RFC 5326 \u2014 deep-space links",
    "Reliable segments + retransmission",
    "Handles long light-time delays",
]
add_card(
    slide, Inches(0.3), Inches(4.65), Inches(3.0), Inches(2.6),
    "LTP (RFC 5326)",
    ltp_bullets,
    border=ACCENT_CYAN,
)

tcpcl_bullets = [
    "TCP Convergence Layer",
    "RFC 7242 \u2014 Earth segment",
    "Session management over TCP",
    "Reliable terrestrial backhaul",
]
add_card(
    slide, Inches(3.5), Inches(4.65), Inches(3.0), Inches(2.6),
    "TCPCL (RFC 7242)",
    tcpcl_bullets,
    border=GREEN,
)

udpcl_bullets = [
    "UDP Convergence Layer",
    "Optical inter-satellite links",
    "Fragmentation + loss simulation",
    "Low-latency mesh forwarding",
]
add_card(
    slide, Inches(6.7), Inches(4.65), Inches(3.0), Inches(2.6),
    "UDP-CL",
    udpcl_bullets,
    border=ACCENT_ORANGE,
)

add_footer(slide, 9, citations="[9] RFC 9171  ·  [10] RFC 5326  ·  [11] RFC 7242")

# ── SLIDE 14 ── Network Topology ──────────────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "push")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "NETWORK TOPOLOGY"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "241 Nodes, 5 Tiers"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

tiers = [
    ("T1  Earth Ground", "6 nodes", "DSN stations: Goldstone, Madrid, Canberra", ACCENT_BLUE),
    ("T2  Earth Orbital", "51 nodes", "GEO relays + 48 LEO laser constellation", ACCENT_PURPLE),
    ("T3  Deep Space Transit", "4 nodes", "Lagrange point relays (ES-L4, ES-L5)", ACCENT_ORANGE),
    ("T4  Mars Orbital", "4 nodes", "Areostationary + polar orbit relays", ACCENT_RED),
    ("T5  Mars Surface", "176 nodes", "Bases, rovers, drones, sensor networks", YELLOW),
]

tier_y = Inches(1.6)
for tier_name, node_count, desc, color in tiers:
    bar = add_shape(
        slide, Inches(0.4), tier_y, Inches(0.08), Inches(0.65),
        fill_color=color, shape_type=MSO_SHAPE.RECTANGLE,
    )
    row_bg = add_shape(
        slide, Inches(0.48), tier_y, Inches(4.9), Inches(0.65),
        fill_color=CARD_BG, border_color=color, shape_type=MSO_SHAPE.RECTANGLE,
    )
    row_bg.line.color.rgb = color
    row_bg.line.width = Pt(1)

    tf = row_bg.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_top = Pt(2)
    p_name = tf.paragraphs[0]
    p_name.text = tier_name
    p_name.font.size = Pt(11)
    p_name.font.bold = True
    p_name.font.color.rgb = WHITE

    p_nodes = add_paragraph(tf)
    p_nodes.text = f"{node_count} \u2014 {desc}"
    p_nodes.font.size = Pt(8)
    p_nodes.font.color.rgb = LIGHT_GRAY

    tier_y += Inches(0.78)

add_image_safe(slide, os.path.join(DIAGRAMS_DIR, "5tier_network.png"), Inches(5.6), Inches(1.5), Inches(4.2), Inches(3.5))

callout_bg = add_shape(
    slide, Inches(0.4), Inches(5.7), Inches(9.2), Inches(0.65),
    fill_color=CARD_BG, border_color=ACCENT_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
)
callout_bg.line.color.rgb = ACCENT_BLUE
callout_bg.line.width = Pt(1)
ctf = callout_bg.text_frame
ctf.word_wrap = True
ctf.margin_left = Pt(12)
cp = ctf.paragraphs[0]
cp.text = "Multiple redundant paths \u00b7 No single point of failure \u00b7 Lagrange relays for conjunction coverage"
cp.font.size = Pt(11)
cp.font.color.rgb = ACCENT_CYAN
cp.alignment = PP_ALIGN.CENTER

add_footer(slide, 10, citations="[A2] AETHERIX topology.py (241 nodes, 5 tiers)  ·  [1] NASA Deep Space Network")

# ── SLIDE 15 ── 5-Tier Network Diagram ────────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "wipe")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "5-TIER NETWORK DIAGRAM"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5))

tier_headers = ["Tier", "Nodes", "Description"]
tier_rows = [
    ["T1 \u2014 Earth Ground", "6", "DSN stations (Goldstone, Madrid, Canberra)"],
    ["T2 \u2014 Earth Orbital", "51", "3 GEO relays + 48 LEO laser satellites"],
    ["T3 \u2014 Deep Space", "4", "ES-L4, ES-L5 Lagrange point relays"],
    ["T4 \u2014 Mars Orbital", "4", "Areostationary relay + polar orbit relays"],
    ["T5 \u2014 Mars Surface", "176", "Habitats, rovers, drones, sensor networks"],
]
_d1 = [tier_headers] + tier_rows
add_table(
    slide, Inches(0.4), Inches(1.2), Inches(9.2), Inches(2.3), len(_d1), len(tier_headers), _d1,
    header_color=ACCENT_BLUE,
)

link_headers = ["Segment", "Data Rate", "Technology"]
link_rows = [
    ["Earth \u2194 Deep Space", "100 Mbps", "1550 nm optical laser"],
    ["Deep Space \u2194 Mars", "2\u2013200 Mbps", "Optical (distance-dependent)"],
    ["Mars Orbital \u2194 Surface", "2 Mbps", "UHF S-band radio"],
    ["LEO Inter-Satellite", "10 Gbps", "Laser ISL mesh"],
]
_d2 = [link_headers] + link_rows
add_table(
    slide, Inches(0.4), Inches(3.8), Inches(9.2), Inches(1.9), len(_d2), len(link_headers), _d2,
    header_color=ACCENT_BLUE,
)

add_footer(slide, 11, citations="[A2] topology.py  ·  [3] JPL Horizons (Lagrange ES-L4/L5 geometry)")

# ── SLIDE 16 ── Network Diagram Visual ────────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "cover")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "NETWORK DIAGRAM"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5))

add_image_safe(
    slide, os.path.join(DIAGRAMS_DIR, "5tier_network.png"),
    Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.3),
)

add_footer(slide, 12, citations="[A2] topology.py  ·  [3] JPL Horizons (Lagrange ES-L4/L5 geometry)")

# --- SLIDE 17 — Chart: DSN Coverage ---
print("Creating Slide 17: Chart \u2014 DSN Coverage...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "dsn_coverage.png"),
    "DSN COVERAGE", "Deep Space Network Ground Stations",
    "Goldstone, Madrid, and Canberra provide 24/7 coverage with 120\u00b0 spacing", ACCENT_BLUE, citations="[1] NASA Deep Space Network"
)

# --- SLIDE 18 — Chart: Orbital Positions ---
print("Creating Slide 18: Chart \u2014 Orbital Positions...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "orbital_positions.png"),
    "ORBITAL POSITIONS", "Key Relay Positions in the Network",
    "Lagrange point relays at ES-L4 and ES-L5 maintain connectivity during conjunction", ACCENT_ORANGE, citations="[3] JPL Horizons (ES-L4/L5)  ·  [A2] topology.py"
)

# ── SLIDE 19 ── Optical Communications ────────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "OPTICAL COMMUNICATIONS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "10\u2013100x Capability Over RF"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

add_stat_card(slide, Inches(0.3), Inches(1.55), Inches(3.1), Inches(1.2), "100\u2013200", "Mbps at closest approach", GREEN)
add_stat_card(slide, Inches(3.55), Inches(1.55), Inches(3.1), Inches(1.2), "10\u201320", "Mbps average distance", ACCENT_BLUE)
add_stat_card(slide, Inches(6.8), Inches(1.55), Inches(3.1), Inches(1.2), "2\u20135", "Mbps at farthest point", ACCENT_RED)

eq_bullets = [
    "FSPL = 20\u00b7log\u2081\u2080(4\u03c0d/\u03bb)  \u2014 free-space path loss",
    "Gain = 10\u00b7log\u2081\u2080(\u03b7\u00b7(\u03c0D/\u03bb)\u00b2)  \u2014 antenna gain",
    "Pr = Pt + Gt + Gr \u2212 FSPL  \u2014 received power",
]
add_card(
    slide, Inches(0.3), Inches(2.9), Inches(4.8), Inches(1.8),
    "Key Equations",
    eq_bullets,
    border=ACCENT_BLUE,
)

config_headers = ["Parameter", "Value"]
config_rows = [
    ["Wavelength", "1550 nm"],
    ["Tx Power", "5 W"],
    ["Tx Aperture", "22 cm"],
    ["Rx Aperture", "1.0 m"],
]
_d3 = [config_headers] + config_rows
add_table(
    slide, Inches(5.3), Inches(2.9), Inches(4.4), Inches(1.8), len(_d3), len(config_headers), _d3,
    header_color=ACCENT_BLUE,
)

add_footer(slide, 13, citations="[5] CCSDS 141.0-B-1 (optical link)  ·  [A4] AETHERIX link_budget.py  ·  [1] NASA MRO (6 Mbps RF baseline)")

# --- SLIDE 20 — Chart: Data Rate vs Distance ---
print("Creating Slide 20: Chart \u2014 Data Rate vs Distance...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "data_rate_vs_distance.png"),
    "DATA RATE VS DISTANCE", "Optical Link Performance",
    "Data rate degrades with distance squared \u2014 200 Mbps at closest to 2 Mbps at farthest", ACCENT_BLUE, citations="[5] CCSDS 141.0-B-1  ·  [A4] link_budget.py  ·  design target",
)

# --- SLIDE 21 — Chart: Link Budget Breakdown ---
print("Creating Slide 21: Chart \u2014 Link Budget Breakdown...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "link_budget_breakdown.png"),
    "LINK BUDGET BREAKDOWN", "Optical Link Analysis",
    "Free-space path loss is the dominant factor \u2014 over 350 dB at average distance", ACCENT_CYAN, citations="[5] CCSDS 141.0-B-1  ·  [A4] AETHERIX link_budget.py"
)

# ── SLIDE 22 ── Earth-Mars Journey ────────────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "push")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "EARTH-MARS JOURNEY"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "500 MB in 7 Hops"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

hop_headers = ["Hop", "Path", "Link"]
hop_rows = [
    ["1", "Rover \u2192 UHF Relay", "UHF S-band"],
    ["2", "UHF Relay \u2192 Areostationary", "UHF uplink"],
    ["3", "Areostationary \u2192 Polar Orbiter", "Crosslink"],
    ["4\u20135", "Polar Orbiter \u2192 LEO Constellation", "1550 nm laser"],
    ["6", "LEO Mesh \u2192 DSN Ground Station", "Optical downlink"],
    ["7", "DSN \u2192 Mission Operations Center", "TCP/IP fiber"],
]
_d4 = [hop_headers] + hop_rows
add_table(
    slide, Inches(0.3), Inches(1.5), Inches(9.4), Inches(2.7), len(_d4), len(hop_headers), _d4,
    header_color=ACCENT_BLUE,
)

add_stat_card(slide, Inches(0.2), Inches(4.4), Inches(2.35), Inches(1.1), "~13min", "End-to-end latency", ACCENT_BLUE)
add_stat_card(slide, Inches(2.65), Inches(4.4), Inches(2.35), Inches(1.1), "<5%", "Protocol overhead", GREEN)
add_stat_card(slide, Inches(5.1), Inches(4.4), Inches(2.35), Inches(1.1), "Tgt >98%", "Delivery (goal)", ACCENT_PURPLE)
add_stat_card(slide, Inches(7.55), Inches(4.4), Inches(2.35), Inches(1.1), "7 hops", "Store-and-forward", ACCENT_ORANGE)

add_footer(slide, 14, citations="[A2] topology.py (7-hop path)  ·  [3] JPL Horizons (light-time)  ·  delivery/throughput are design targets")

# --- SLIDE 23 — Chart: Latency Comparison ---
print("Creating Slide 23: Chart \u2014 Latency Comparison...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "latency_comparison.png"),
    "LATENCY COMPARISON", "Protocol Performance",
    "DTN overhead adds less than 5% to the pure light-time delay", ACCENT_BLUE, citations="[9][10][11] IETF RFCs  ·  [12] RFC 4838  ·  [A3] run_simulation"
)

# --- SLIDE 24 — Chart: Data Volume ---
print("Creating Slide 24: Chart \u2014 Data Volume...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "data_volume.png"),
    "DATA VOLUME ANALYSIS", "Daily Throughput Comparison",
    "AETHERIX delivers 10\u201320x more data per day than current systems", GREEN, citations="[A2] topology.py  ·  [1] NASA MRO baseline  ·  design targets"
)

# ── SLIDE 25 ── RL Routing ────────────────────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "wipe")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "RL-BASED ROUTING"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "AI Innovation"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

cgr_bullets = [
    "Pre-computed contact plans required",
    "Cannot adapt to unexpected link failures",
    "Computational cost grows O(n\u00b3) per route",
    "No learning from historical performance",
    "Single-path \u2014 ignores multipath opportunity",
]
add_card(
    slide, Inches(0.3), Inches(1.55), Inches(4.6), Inches(2.5),
    "Contact Graph Routing",
    cgr_bullets,
    border=ACCENT_RED,
)

rl_bullets = [
    "State: (node, neighbors, link quality, buffer)",
    "Actions: forward | store | drop | split",
    "Reward: R = \u03b1(delivery) \u2212 \u03b2(delay) \u2212 \u03b3(hops) \u2212 \u03b4(drops) \u2212 \u03b5(energy)",
    "\u03b5-greedy exploration with decay",
    "Q-table + federated aggregation across nodes",
]
add_card(
    slide, Inches(5.1), Inches(1.55), Inches(4.6), Inches(2.5),
    "RL Agent",
    rl_bullets,
    border=ACCENT_CYAN,
)

add_stat_card(slide, Inches(0.3), Inches(4.25), Inches(3.1), Inches(1.1), "Tgt +20\u201340%", "vs CGR (goal)", GREEN)
add_stat_card(slide, Inches(3.55), Inches(4.25), Inches(3.1), Inches(1.1), "Tgt 1000\u00d7", "Recovery (goal)", ACCENT_BLUE)
add_stat_card(slide, Inches(6.8), Inches(4.25), Inches(3.1), Inches(1.1), "Federated", "Multi-agent learning", ACCENT_PURPLE)

add_footer(slide, 15, citations="[A1] rl_agent.py (reward fn, ε-decay 0.995)  ·  [A3] run_simulation Module 3 (training convergence 713/800)")

# --- SLIDE 26 — Chart: RL Routing Heatmap ---
print("Creating Slide 26: Chart \u2014 RL Routing Heatmap...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "rl_routing_heatmap.png"),
    "RL ROUTING HEATMAP", "Q-Value Distribution",
    "Learned routing preferences across network nodes \u2014 brighter = higher Q-value", ACCENT_CYAN, citations="[A1] rl_agent.py (Q-table, ε-greedy)  ·  [A3] Module 3"
)

# --- SLIDE 27 — Chart: Energy Efficiency ---
print("Creating Slide 27: Chart \u2014 Energy Efficiency...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "energy_efficiency.png"),
    "ENERGY EFFICIENCY", "RL Agent Power Optimization",
    "The reward function penalizes energy waste, optimizing for mission lifetime", ACCENT_PURPLE, citations="[A1] rl_agent.py (energy penalty in reward fn)"
)

# ── SLIDE 28 ── Quantum Security ──────────────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "cover")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "QUANTUM SECURITY"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Future-Proof"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

bb84_bullets = [
    "1. Alice sends random polarized photons (H/V, D/A bases)",
    "2. Bob measures each photon in a random basis",
    "3. Public reconciliation \u2014 keep matching-basis bits",
    "4. Error estimation \u2014 sample subset for QBER",
    "5. Privacy amplification \u2014 universal hashing",
    "6. QBER < 11% \u2192 SECURE key established",
]
add_card(
    slide, Inches(0.3), Inches(1.55), Inches(4.8), Inches(2.5),
    "BB84 Protocol",
    bb84_bullets,
    border=ACCENT_PURPLE,
)

pqc_bullets = [
    "Kyber (ML-KEM) \u2014 key encapsulation",
    "Dilithium (ML-DSA) \u2014 digital signatures",
    "NIST PQC standard \u2014 quantum-resistant",
    "Hybrid classical-quantum key exchange",
]
add_card(
    slide, Inches(5.3), Inches(1.55), Inches(4.4), Inches(2.5),
    "Post-Quantum Cryptography",
    pqc_bullets,
    border=ACCENT_ORANGE,
)

deploy_headers = ["Phase", "Protocol", "Key Rate"]
deploy_rows = [
    ["Phase 1 \u2014 LEO", "BB84", "1\u201310 kbps"],
    ["Phase 2 \u2014 GEO", "BB84 + E91", "10\u2013100 kbps"],
    ["Phase 3 \u2014 Mars", "E91 + Repeaters", "1+ kbps"],
]
_d5 = [deploy_headers] + deploy_rows
add_table(
    slide, Inches(0.3), Inches(4.25), Inches(9.4), Inches(1.5), len(_d5), len(deploy_headers), _d5,
    header_color=ACCENT_PURPLE,
)

add_footer(slide, 16, citations="[13] Bennett-Brassard 1984  ·  [14] Ekert 1991  ·  [15] Shor-Preskill 2000 (QBER<11%)  ·  [16][17] NIST FIPS 203/204  ·  [A5] qkd.py")

# --- SLIDE 29 — Chart: QKD Security ---
print("Creating Slide 29: Chart \u2014 QKD Security...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "qkd_security.png"),
    "QKD SECURITY ANALYSIS", "Quantum Bit Error Rate",
    "QBER below 11% confirms secure key exchange \u2014 no eavesdropper detected", ACCENT_PURPLE, citations="[13] Bennett-Brassard 1984  ·  [15] Shor-Preskill 2000 (QBER<11%)  ·  [A5] qkd.py"
)

# --- SLIDE 30 — Chart: QKD Key Rate ---
print("Creating Slide 30: Chart \u2014 QKD Key Rate...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "qkd_key_rate.png"),
    "QKD KEY GENERATION RATE", "Secure Key Rate vs Distance",
    "Key rate decreases with channel loss \u2014 repeater chains extend operational range", ACCENT_CYAN, citations="[13] BB84  ·  [14] Ekert 1991  ·  [A5] qkd.py + repeater_chain.py"
)

# ── SLIDE 31 ── Orbital Mechanics ─────────────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "ORBITAL MECHANICS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Contact Windows"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.LEFT

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5))

mars_headers = ["Parameter", "Value"]
mars_rows = [
    ["Semi-major axis", "1.524 AU"],
    ["Synodic period", "779.94 days"],
    ["Distance range", "54.6M \u2013 401M km"],
    ["Areostationary altitude", "17,032 km"],
    ["One-way light time", "3 \u2013 22 minutes"],
]
_d6 = [mars_headers] + mars_rows
add_table(
    slide, Inches(0.3), Inches(1.55), Inches(4.7), Inches(2.3), len(_d6), len(mars_headers), _d6,
    header_color=ACCENT_ORANGE,
)

window_headers = ["Window", "Availability", "Duration", "Data Rate"]
window_rows = [
    ["Optimal", "99%", "8\u201312 hrs", "100\u2013200 Mbps"],
    ["Good", "95%", "6\u20138 hrs", "20\u2013100 Mbps"],
    ["Fair", "85%", "4\u20136 hrs", "5\u201320 Mbps"],
    ["Blackout", "0%", "2\u20134 weeks", "\u2014"],
]
_d7 = [window_headers] + window_rows
add_table(
    slide, Inches(5.1), Inches(1.55), Inches(4.6), Inches(2.3), len(_d7), len(window_headers), _d7,
    header_color=ACCENT_BLUE,
)

add_footer(slide, 17, citations="[3] JPL Horizons (synodic 779.94 d, 54.6M–401M km)  ·  [A2] topology.py  ·  [A4] link_budget.py")

# --- SLIDE 32 — Chart: Contact Windows ---
print("Creating Slide 32: Chart \u2014 Contact Windows...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "contact_windows.png"),
    "CONTACT WINDOWS", "Communication Opportunity Prediction",
    "Contact windows vary from 4 to 12 hours depending on orbital geometry", ACCENT_ORANGE, citations="[3] JPL Horizons  ·  [A2] topology.py / contact_windows.py"
)

# ── SLIDE 33 ── Radiation-Hardened Computing ──────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "RADIATION-HARDENED COMPUTING"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(30), True, WHITE, PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Surviving SEUs, latchup and total dose on the way to Mars"
p2.font.size, p2.font.color.rgb = Pt(15), ACCENT_CYAN

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5), ACCENT_RED)

effect_data = [
    ["Effect", "What it does", "Mitigation"],
    ["SEU", "Single bit flip", "SECDED ECC"],
    ["MBU", "Multi-bit flip (1 ion)", "Bit interleaving"],
    ["SEL", "Latchup (destructive)", "Current limit + power-cycle"],
    ["TID", "Cumulative dose", "Rad-hard parts (RAD750)"],
]
add_table(slide, Inches(0.3), Inches(1.6), Inches(6.4), Inches(2.4),
          len(effect_data), len(effect_data[0]), effect_data, header_color=ACCENT_RED)

mitig_card = add_card(
    slide, Inches(7.0), Inches(1.6), Inches(6.0), Inches(2.4),
    title="Defense-in-Depth Stack",
    body_lines=[
        "\u2022 TMR \u2014 triple replicas, majority vote (masks logic faults)",
        "\u2022 SECDED (39,32) ECC \u2014 correct 1 bit, detect 2",
        "\u2022 Scrubbing \u2014 rewrite memory before a 2nd upset accumulates",
        "\u2022 FDIR + watchdog \u2014 detect \u2192 isolate \u2192 reset \u2192 SAFE-MODE",
    ],
    title_color=ACCENT_RED,
)

add_stat_card(slide, Inches(0.3), Inches(4.4), Inches(3.0), Inches(1.2),
              "200x", "Fewer errors (ECC + scrub + interleave)", value_color=ACCENT_CYAN)
add_stat_card(slide, Inches(3.5), Inches(4.4), Inches(3.0), Inches(1.2),
              "3,334x", "TMR reliability gain (p=1e-4/op)", value_color=ACCENT_BLUE)
add_stat_card(slide, Inches(6.7), Inches(4.4), Inches(3.0), Inches(1.2),
              "200 krad", "RAD750 TID tolerance (>2000x margin)", value_color=ACCENT_ORANGE)
add_stat_card(slide, Inches(9.9), Inches(4.4), Inches(3.1), Inches(1.2),
              "~0.9 / day", "Residual uncorrectable, Earth-Mars transit", value_color=ACCENT_PURPLE)

note = add_textbox(slide, Inches(0.6), Inches(5.85), Inches(12.2), Inches(0.9))
pn = note.text_frame.paragraphs[0]
pn.text = ("Model: 512 Mbit, ~210-day GCR cruise. ~37,000 raw bit upsets reduced to "
           "~186 uncorrectable over the mission. Heritage: NASA RAD750 (Curiosity/"
           "Perseverance), ESA LEON3FT.  \u2192  src/computing/radiation.py")
pn.font.size, pn.font.color.rgb = Pt(11), MED_GRAY
add_footer(slide, citations="[A6] AETHERIX radiation.py (Module 6: 37,159 upsets, 200×/3,334×/2,127×)  ·  [18] BAE RAD750  ·  [19] ESA LEON3FT")

# ── SLIDE 34 ── Mission-Critical Data Prioritization ──────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "push")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "MISSION-CRITICAL DATA PRIORITIZATION"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Bandwidth triage: get the right bits home first"
p2.font.size, p2.font.color.rgb = Pt(15), ACCENT_CYAN

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5), ACCENT_ORANGE)

tier_data = [
    ["Tier", "Class", "Examples"],
    ["P0", "Emergency / Safety", "Health telemetry, collision avoidance"],
    ["P1", "Mission-critical", "Command ACKs, time-sensitive science"],
    ["P2", "High-priority", "Routine telemetry, scheduled science"],
    ["P4", "Low / Bulk", "Housekeeping logs, file transfers"],
]
add_table(slide, Inches(0.3), Inches(1.6), Inches(7.0), Inches(2.4),
          len(tier_data), len(tier_data[0]), tier_data, header_color=ACCENT_ORANGE)

comp_data = [
    ["Data type", "Standard", "Ratio"],
    ["Telemetry", "CCSDS 121", "3x"],
    ["Imagery (lossy)", "CCSDS 122", "10x"],
    ["Video", "H.265", "50x"],
]
add_table(slide, Inches(7.6), Inches(1.6), Inches(5.4), Inches(2.0),
          len(comp_data), len(comp_data[0]), comp_data, header_color=ACCENT_PURPLE)

add_stat_card(slide, Inches(0.3), Inches(4.4), Inches(3.0), Inches(1.2),
              "100%", "Link utilization (target)", value_color=ACCENT_CYAN)
add_stat_card(slide, Inches(3.5), Inches(4.4), Inches(3.0), Inches(1.2),
              "5 / 6", "Items fully delivered by priority", value_color=GREEN)
add_stat_card(slide, Inches(6.7), Inches(4.4), Inches(3.0), Inches(1.2),
              "BPv7", "Fragmentation defers bulk remainder", value_color=ACCENT_BLUE)
add_stat_card(slide, Inches(9.9), Inches(4.4), Inches(3.1), Inches(1.2),
              "Preempt", "Emergency uses direct-to-Earth backup", value_color=ACCENT_RED)

note = add_textbox(slide, Inches(0.6), Inches(5.85), Inches(12.2), Inches(0.9))
pn = note.text_frame.paragraphs[0]
pn.text = ("Scenario: 30 Mbps, 15-min contact, oversubscribed. Deadline-aware, "
           "preemptive QoS scheduler delivers emergency + mission + science first; "
           "6 GB software update fragmented to the next pass.  \u2192  "
           "src/routing/prioritization.py")
pn.font.size, pn.font.color.rgb = Pt(11), MED_GRAY
add_footer(slide, citations="[A7] AETHERIX prioritization.py  ·  [7] CCSDS 121.0-B-3  ·  [8] CCSDS 122.0-B-2  ·  [9] RFC 9171")

# ── SLIDE ── Engineering Trade-offs ───────────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "ENGINEERING TRADE-OFFS"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(30), True, WHITE, PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Every design choice is a trade-off \u2014 here are the four we made"
p2.font.size, p2.font.color.rgb = Pt(15), ACCENT_ORANGE

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5), ACCENT_ORANGE)

trade_data = [
    ["Decision", "Chose", "Trade-off accepted"],
    ["Optical vs RF", "1550 nm laser", "Needs clear line-of-sight; clouds block"],
    ["DTN vs TCP", "Store-and-forward", "Storage + custody overhead (~5%)"],
    ["RL vs CGR", "Q-learning adaptive", "Training time; exploration risk early on"],
    ["QKD vs AES", "BB84/E91 + PQC", "Key rate drops with distance; HW complexity"],
]
add_table(slide, Inches(0.3), Inches(1.6), Inches(12.6), Inches(2.8),
          len(trade_data), len(trade_data[0]), trade_data, header_color=ACCENT_ORANGE)

trade_card = add_card(
    slide, Inches(0.6), Inches(4.7), Inches(12.0), Inches(1.8),
    title="Net result of the trade-offs",
    body_lines=[
        "\u2022 10\u2013100\u00d7 more data than RF-only systems  \u2192  worth the line-of-sight constraint",
        "\u2022 <5% DTN overhead buys tolerance to 22-minute disruptions  \u2192  worth the storage cost",
        "\u2022 RL adds autonomy (no human replanning)  \u2192  worth the convergence period",
        "\u2022 QKD + PQC gives information-theoretic security  \u2192  worth the key-rate limit",
    ],
    title_color=ACCENT_ORANGE,
)

add_footer(slide, citations="[A4] link_budget.py (optical vs RF)  ·  [A1] rl_agent.py (RL vs CGR)  ·  [A5] qkd.py (QKD vs AES)")

# ── SLIDE ── Failure Detection & Recovery ─────────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "FAILURE DETECTION & RECOVERY"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(30), True, WHITE, PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "When links fail, the RL agent reroutes through Lagrange relays"
p2.font.size, p2.font.color.rgb = Pt(15), ACCENT_CYAN

add_accent_line(slide, Inches(0.6), Inches(1.35), Inches(2.5), ACCENT_CYAN)

fail_data = [
    ["Scenario", "Route", "Reward", "Outcome"],
    ["Optical link OPEN", "Direct GEO \u2192 Areostation", "\u20131.438", "Baseline path"],
    ["ES-L4/L5 OPEN", "Reroute via Lagrange relay", "\u20130.201", "7\u00d7 better"],
    ["All links up", "RL optimal path", "+0.156", "Full throughput"],
]
add_table(slide, Inches(0.3), Inches(1.6), Inches(12.6), Inches(2.4),
          len(fail_data), len(fail_data[0]), fail_data, header_color=ACCENT_CYAN)

recovery_card = add_card(
    slide, Inches(0.6), Inches(4.3), Inches(6.0), Inches(2.4),
    title="How recovery works",
    body_lines=[
        "\u2022 RL agent detects link failure via Q-value collapse",
        "\u2022 Explores alternate neighbours (ES-L4, ES-L5, LEO mesh)",
        "\u2022 Converges on best alternate in <1 epoch",
        "\u2022 Custody transfer guarantees no bundle loss",
    ],
    title_color=ACCENT_CYAN,
)

resilience_card = add_card(
    slide, Inches(6.9), Inches(4.3), Inches(6.0), Inches(2.4),
    title="Resilience metrics",
    body_lines=[
        "\u2022 Solar conjunction: 50\u201370% throughput via L4/L5 vs 0% blackout",
        "\u2022 Node failure: <1 epoch rerouting (federated Q-table)",
        "\u2022 Bundle loss: 0 (custody + LTP retransmission)",
        "\u2022 Mean time to recover: dominated by light-time, not compute",
    ],
    title_color=ACCENT_CYAN,
)

add_footer(slide, citations="[A8] AETHERIX run_simulation Module 4 (\u20131.438 \u2192 \u20130.201 demonstrated)  ·  [A1] rl_agent.py  ·  [3] JPL Horizons (Lagrange geometry)")

# ── SLIDE 35 ── End-to-End Mission ────────────────────────────────────────────

slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "push")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "END-TO-END MISSION"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Mars Surface \u2192 Earth"
p2.font.size, p2.font.color.rgb, p2.alignment = Pt(14), ACCENT_CYAN, PP_ALIGN.LEFT

hops = [
    "Rover", "UHF", "Areostationary", "Polar",
    "Deep Space 1550 nm", "LEO", "DSN / MOC",
]
hop_x_start = 0.4
hop_spacing = (SLIDE_WIDTH - Inches(2 * hop_x_start)) / len(hops)
for i, hop in enumerate(hops):
    cx = Inches(hop_x_start) + int(hop_spacing * i) + int(hop_spacing / 2) - Inches(0.45)
    cy = Inches(1.7)
    shape = add_shape(
        slide, Inches(0.4), Inches(0.4), Inches(0.9), Inches(0.45),
        fill_color=ACCENT_BLUE, border_color=CARD_BORDER, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    tf = shape.text_frame
    tf.word_wrap = True
    pp = tf.paragraphs[0]
    pp.text = hop
    pp.font.size, pp.font.bold, pp.font.color.rgb, pp.alignment = Pt(9), True, WHITE, PP_ALIGN.CENTER
    if i < len(hops) - 1:
        arrow_left = int(cx) + Inches(0.9)
        arrow_top = int(cy) + Inches(0.15)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, Inches(0.2), Inches(0.15))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_CYAN
        arrow.line.fill.background()

stats_18 = [
    ("~13 min", "End-to-End Latency", ACCENT_BLUE),
    ("<5%", "Packet Loss", GREEN),
    ("Tgt >98%", "Delivery (goal)", ACCENT_PURPLE),
    ("7 hops", "Network Path", ACCENT_ORANGE),
]
stat_w = Inches(2.1)
stat_h = Inches(1.1)
stat_gap = Inches(0.25)
total_w = 4 * stat_w + 3 * stat_gap
stat_x0 = (SLIDE_WIDTH - total_w) // 2
for idx, (val, label, color) in enumerate(stats_18):
    sx = stat_x0 + int((stat_w + stat_gap) * idx)
    add_stat_card(slide, sx, Inches(2.55), stat_w, stat_h, val, label, color)

fail_x = Inches(0.6)
fail_y = Inches(4.0)
fail_w = SLIDE_WIDTH - Inches(1.2)
fail_h = Inches(1.2)
fail_card = add_card(slide, fail_x, fail_y, fail_w, fail_h)
fail_card.line.color.rgb = ACCENT_RED
fail_card.line.width = Pt(2.5)
tf_fail = fail_card.text_frame
tf_fail.word_wrap = True
p_fail_title = tf_fail.paragraphs[0]
p_fail_title.text = "\u26a0 FAILURE SCENARIO"
p_fail_title.font.size, p_fail_title.font.bold, p_fail_title.font.color.rgb = Pt(12), True, ACCENT_RED
p_fail_body = tf_fail.add_paragraph()
p_fail_body.text = (
    "Solar flare at T+400s \u2192 stored at MRS-Polar \u2192 "
    "rerouted via ES-L4 \u2192 delayed ~30 min NOT LOST"
)
p_fail_body.font.size, p_fail_body.font.color.rgb = Pt(11), WHITE

add_footer(slide, 18, citations="[A2] topology.py  ·  [A3] run_simulation  ·  targets clearly labelled")

# ── SLIDE 36 ── Data Flow Diagram (Text) ────────────────────────────────────────

slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "wipe")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "DATA FLOW DIAGRAM"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

card_w = Inches(4.3)
card_h = Inches(3.8)
card_y = Inches(1.3)

left_card = add_card(slide, Inches(0.5), card_y, card_w, card_h)
tf_l = left_card.text_frame
tf_l.word_wrap = True
p_title_l = tf_l.paragraphs[0]
p_title_l.text = "Application \u2192 Transport"
p_title_l.font.size, p_title_l.font.bold, p_title_l.font.color.rgb = Pt(14), True, ACCENT_CYAN
p_title_l.space_after = Pt(10)

left_steps = [
    ("1. Source Node", "Science data generated"),
    ("2. Bundle Protocol", "BPv7 wraps data + metadata"),
    ("3. RL Routing", "Agent evaluates state, selects hop"),
    ("4. QKD Encrypt", "BB84 shared key applied"),
]
for step_title, step_desc in left_steps:
    p_t = tf_l.add_paragraph()
    p_t.text = step_title
    p_t.font.size, p_t.font.bold, p_t.font.color.rgb = Pt(11), True, WHITE
    p_t.space_before = Pt(6)
    p_d = tf_l.add_paragraph()
    p_d.text = step_desc
    p_d.font.size, p_d.font.color.rgb = Pt(9), LIGHT_GRAY

right_card = add_card(slide, Inches(5.2), card_y, card_w, card_h)
tf_r = right_card.text_frame
tf_r.word_wrap = True
p_title_r = tf_r.paragraphs[0]
p_title_r.text = "Convergence \u2192 Physical \u2192 Delivery"
p_title_r.font.size, p_title_r.font.bold, p_title_r.font.color.rgb = Pt(14), True, ACCENT_CYAN
p_title_r.space_after = Pt(10)

right_steps = [
    ("1. LTP Segmentation", "Bundle split into blocks"),
    ("2. Store & Wait", "Buffer until link available"),
    ("3. Physical TX", "UHF \u2192 Optical ISL \u2192 1550 nm \u2192 Ka-band"),
    ("4. LTP Reassemble", "Reassemble \u2192 Decrypt \u2192 Deliver"),
]
for step_title, step_desc in right_steps:
    p_t = tf_r.add_paragraph()
    p_t.text = step_title
    p_t.font.size, p_t.font.bold, p_t.font.color.rgb = Pt(11), True, WHITE
    p_t.space_before = Pt(6)
    p_d = tf_r.add_paragraph()
    p_d.text = step_desc
    p_d.font.size, p_d.font.color.rgb = Pt(9), LIGHT_GRAY

add_footer(slide, 19, citations="[A2] topology.py (full data path Mars→Earth)")

# ── SLIDE 37 ── Data Flow Visual ────────────────────────────────────────────────

slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "cover")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "DATA FLOW DIAGRAM"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

img_path_20 = os.path.join(DIAGRAMS_DIR, "data_flow.png")
add_image_safe(
    slide, img_path_20,
    Inches(0.4), Inches(1.2), Inches(9.2), Inches(5.5),
)

add_footer(slide, 20, citations="[A2] topology.py (full data path Mars→Earth)")

# ── SLIDE 38 ── Performance ────────────────────────────────────────────────────

slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "PERFORMANCE COMPARISON"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(9), Inches(0.35))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "AETHERIX vs Current"
p2.font.size, p2.font.color.rgb, p2.alignment = Pt(14), ACCENT_CYAN, PP_ALIGN.LEFT

perf_data = [
    ["Metric", "Current (MRO)", "AETHERIX", "Improvement"],
    ["Downlink Rate", "0.5-6 Mbps", "2-200 Mbps", "10-100\u00d7 capability"],
    ["Daily Volume", "5-10 GB", "50-100 GB", "10-20\u00d7 (target)"],
    ["Availability", "60-75%", ">95% (tgt)", "+20-35%"],
    ["Routing", "Static (CGR)", "RL-adaptive", "Autonomous"],
    ["Security", "AES-256", "QKD + PQC", "Quantum-proof"],
    ["Scalability", "5-10 assets", "241 nodes", "24-48\u00d7"],
    ["Conjunction", "Blackout", "L4/L5 relay (target)", "+50-70% (target)"],
]
add_table(
    slide, Inches(0.4), Inches(1.25), Inches(9.2), Inches(3.4), len(perf_data), len(perf_data[0]) if perf_data else 0, perf_data,
    header_color=GREEN,
)

add_footer(slide, 21, citations="[1] NASA MRO (0.5–6 Mbps)  ·  [A4] link_budget.py (2–200 Mbps capability)  ·  [A2] topology.py  ·  [A8] Module 4  ·  targets clearly labelled")

# --- SLIDE 39 — Chart: Performance Comparison ---
print("Creating Slide 39: Chart \u2014 Performance Comparison...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "performance_comparison.png"),
    "PERFORMANCE COMPARISON", "AETHERIX vs Current Systems",
    "Model-based targets vs current MRO baseline", GREEN, citations="[1] NASA MRO  ·  [A4] link_budget.py  ·  [A2] topology.py  ·  design targets"
)

# --- SLIDE 40 — Chart: Optical vs RF Radar ---
print("Creating Slide 40: Chart \u2014 Optical vs RF Radar...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "optical_vs_rf_radar.png"),
    "OPTICAL vs RF", "Link Technology Comparison",
    "Optical links provide dramatically higher data rates, especially at shorter ranges", ACCENT_BLUE, citations="[5] CCSDS 141.0-B-1 (optical)  ·  [1] NASA MRO (RF)  ·  [A4] link_budget.py"
)

# ── SLIDE 41 ── Implementation ──────────────────────────────────────────────────

slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "push")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "IMPLEMENTATION"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "What We Built"
p2.font.size, p2.font.color.rgb, p2.alignment = Pt(14), ACCENT_CYAN, PP_ALIGN.LEFT

impl_stats = [
    ("25", "Modules", GREEN),
    ("189", "Tests", ACCENT_BLUE),
    ("12", "Demos", ACCENT_PURPLE),
    ("5", "Policies", ACCENT_ORANGE),
]
stat_w = Inches(2.1)
stat_h = Inches(0.95)
stat_gap = Inches(0.25)
total_w_impl = 4 * stat_w + 3 * stat_gap
stat_x0_impl = (SLIDE_WIDTH - total_w_impl) // 2
for idx, (val, label, color) in enumerate(impl_stats):
    sx = stat_x0_impl + int((stat_w + stat_gap) * idx)
    add_stat_card(slide, sx, Inches(1.55), stat_w, stat_h, val, label, color)

modules_card = add_card(slide, Inches(0.5), Inches(2.8), Inches(4.3), Inches(1.7))
tf_mod = modules_card.text_frame
tf_mod.word_wrap = True
p_mod_title = tf_mod.paragraphs[0]
p_mod_title.text = "Core Modules"
p_mod_title.font.size, p_mod_title.font.bold, p_mod_title.font.color.rgb = Pt(13), True, ACCENT_CYAN
p_mod_title.space_after = Pt(6)
for mod_name in ["Link Budget", "RL Agent", "QKD", "Bundle Protocol", "Topology", "Simulation"]:
    p_mod = tf_mod.add_paragraph()
    p_mod.text = f"\u2022  {mod_name}"
    p_mod.font.size, p_mod.font.color.rgb = Pt(11), WHITE

std_card = add_card(slide, Inches(5.1), Inches(2.8), Inches(4.6), Inches(1.7))
tf_std = std_card.text_frame
tf_std.word_wrap = True
p_std_title = tf_std.paragraphs[0]
p_std_title.text = "Standards Compliance"
p_std_title.font.size, p_std_title.font.bold, p_std_title.font.color.rgb = Pt(13), True, ACCENT_CYAN
p_std_title.space_after = Pt(6)
standards_data = [
    ["Standard", "Title"],
    ["RFC 9171", "Bundle Protocol v7"],
    ["RFC 4838", "DTN Architecture"],
    ["RFC 5326", "LTP (deep space)"],
    ["RFC 9172", "BPSec (security)"],
    ["CCSDS 734.2-B-1", "CCSDS Bundle Protocol"],
    ["CCSDS 734.3-B-1", "SABR (contact routing)"],
    ["CCSDS 141.0-B-1", "Optical comms phys."],
]
std_tbl = add_table(
    slide, Inches(5.2), Inches(3.35), Inches(4.4), Inches(1.1), len(standards_data), len(standards_data[0]) if standards_data else 0, standards_data,
    header_color=ACCENT_BLUE,
)

add_footer(slide, 22, citations="[9][10][11][12] IETF RFCs  ·  [2][5] CCSDS standards  ·  [A1]–[A8] AETHERIX modules")

# --- SLIDE 42 — Chart: Bandwidth Evolution ---
print("Creating Slide 42: Chart \u2014 Bandwidth Evolution...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "bandwidth_evolution.png"),
    "BANDWIDTH EVOLUTION", "Deep-Space Communications Timeline",
    "Optical communications represent the next leap in deep-space data rates", ACCENT_CYAN, citations="[1] NASA (historical)  ·  [4] NASA DSOC/Psyche 2023"
)

# --- SLIDE 43 — Chart: Mission Timeline ---
print("Creating Slide 43: Chart \u2014 Mission Timeline...")
add_chart_slide(
    os.path.join(CHARTS_DIR, "mission_timeline.png"),
    "MISSION TIMELINE", "Deployment Roadmap",
    "Phases 1-4 are complete \u2014 Phase 5 targets ns-3 integration", GREEN, citations="[A1]–[A8] AETHERIX modules  ·  [4] NASA DSOC heritage"
)

# ── SLIDE 44 ── Roadmap ────────────────────────────────────────────────────────

slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "wipe")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "ROADMAP"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "From Demo to Deployment"
p2.font.size, p2.font.color.rgb, p2.alignment = Pt(14), ACCENT_CYAN, PP_ALIGN.LEFT

phases = [
    ("Phase 1-4  \u2713 Complete", "Topology \u2022 RL \u2022 QKD \u2022 Web UI \u2022 480 tests", GREEN),
    ("Phase 5  Network Sim", "ns-3 integration, realistic channel models", ACCENT_BLUE),
    ("Phase 6  Production", "DQN agent, ION-DTN stack, hardware-in-loop", ACCENT_PURPLE),
    ("Phase 7  Hardware", "SDR prototyping, optical ground station demo", ACCENT_ORANGE),
]
phase_w = Inches(4.6)
phase_h = Inches(0.7)
phase_gap = Inches(0.15)
phase_x = Inches(0.5)
for idx, (title, desc, color) in enumerate(phases):
    py = Inches(1.55) + int((phase_h + phase_gap) * idx)
    card = add_card(slide, phase_x, py, phase_w, phase_h)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    accent_bar = add_shape(
        slide, Inches(0.4), Inches(0.4), Inches(0.08), Inches(0.2), shape_type=MSO_SHAPE.RECTANGLE,
        fill_color=color,
    )
    tf_ph = card.text_frame
    tf_ph.word_wrap = True
    p_ph = tf_ph.paragraphs[0]
    p_ph.text = title
    p_ph.font.size, p_ph.font.bold, p_ph.font.color.rgb = Pt(12), True, color
    p_ph2 = tf_ph.add_paragraph()
    p_ph2.text = desc
    p_ph2.font.size, p_ph2.font.color.rgb = Pt(9), LIGHT_GRAY

add_footer(slide, 23, citations="[A1]–[A8] AETHERIX source modules  ·  [4] NASA DSOC (Psyche) heritage  ·  DQN/ns-3/ION-DTN on production roadmap")

# ── SLIDE 45 ── Conclusion ──────────────────────────────────────────────────────

slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "cover")

txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "CONCLUSION"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(9), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "AETHERIX Delivers"
p2.font.size, p2.font.color.rgb, p2.alignment = Pt(14), ACCENT_CYAN, PP_ALIGN.LEFT

concl_card_w = Inches(4.3)
concl_card_h = Inches(1.7)

problem_card = add_card(slide, Inches(0.5), Inches(1.6), concl_card_w, concl_card_h)
problem_card.line.color.rgb = ACCENT_RED
problem_card.line.width = Pt(2.5)
tf_prob = problem_card.text_frame
tf_prob.word_wrap = True
p_prob_t = tf_prob.paragraphs[0]
p_prob_t.text = "Problem Solved"
p_prob_t.font.size, p_prob_t.font.bold, p_prob_t.font.color.rgb = Pt(14), True, ACCENT_RED
p_prob_t.space_after = Pt(8)
p_prob_b = tf_prob.add_paragraph()
p_prob_b.text = "TCP/IP can\u2019t work at interplanetary distances"
p_prob_b.font.size, p_prob_b.font.color.rgb = Pt(12), WHITE

built_card = add_card(slide, Inches(5.2), Inches(1.6), concl_card_w, concl_card_h)
built_card.line.color.rgb = GREEN
built_card.line.width = Pt(2.5)
tf_built = built_card.text_frame
tf_built.word_wrap = True
p_built_t = tf_built.paragraphs[0]
p_built_t.text = "What We Built"
p_built_t.font.size, p_built_t.font.bold, p_built_t.font.color.rgb = Pt(14), True, GREEN
p_built_t.space_after = Pt(8)
p_built_b = tf_built.add_paragraph()
p_built_b.text = "BPv7 \u2022 RL Routing \u2022 QKD \u2022 Optical/RF Hybrid"
p_built_b.font.size, p_built_b.font.color.rgb = Pt(12), WHITE

conc_stats = [
    ("10-100\u00d7", "Target Speedup \u2666", ACCENT_BLUE),
    (">95%", "Availability (target)", GREEN),
    ("RL", "Adaptive Routing \u2713", ACCENT_PURPLE),
    ("QKD", "Quantum Security \u2713", ACCENT_ORANGE),
]
stat_w_c = Inches(2.1)
stat_h_c = Inches(1.1)
stat_gap_c = Inches(0.25)
total_w_c = 4 * stat_w_c + 3 * stat_gap_c
stat_x0_c = (SLIDE_WIDTH - total_w_c) // 2
for idx, (val, label, color) in enumerate(conc_stats):
    sx = stat_x0_c + int((stat_w_c + stat_gap_c) * idx)
    add_stat_card(slide, sx, Inches(3.7), stat_w_c, stat_h_c, val, label, color)

add_footer(slide, 24, citations="[1] NASA  ·  [3] JPL Horizons  ·  [9][12] IETF  ·  [13][15] QKD  ·  [A1][A2][A4][A5][A8] AETHERIX")

# ── SLIDE 46 ── Thank You ──────────────────────────────────────────────────────

slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

top_bar = add_shape(
    slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), shape_type=MSO_SHAPE.RECTANGLE,
    fill_color=ACCENT_BLUE,
)

ty_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.6), SLIDE_WIDTH - Inches(1), Inches(0.9)
)
p_ty = ty_box.text_frame.paragraphs[0]
p_ty.text = "THANK YOU"
p_ty.font.size, p_ty.font.bold, p_ty.font.color.rgb, p_ty.alignment = (
    Pt(64), True, WHITE, PP_ALIGN.CENTER,
)

add_accent_line(slide, Inches(3.6), Inches(1.7), Inches(2.8), ACCENT_CYAN)

q_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(2.0), SLIDE_WIDTH - Inches(1), Inches(0.6)
)
p_q = q_box.text_frame.paragraphs[0]
p_q.text = "Questions?"
p_q.font.size, p_q.font.bold, p_q.font.color.rgb, p_q.alignment = (
    Pt(36), True, ACCENT_CYAN, PP_ALIGN.CENTER,
)

ty_stats = [
    ("10-100\u00d7", "Target Speedup \u2666", ACCENT_BLUE),
    (">95%", "Availability (target)", GREEN),
    ("RL", "Adaptive Routing \u2713", ACCENT_PURPLE),
    ("QKD", "Quantum Security \u2713", ACCENT_ORANGE),
]
stat_w_ty = Inches(2.1)
stat_h_ty = Inches(0.95)
stat_gap_ty = Inches(0.25)
total_w_ty = 4 * stat_w_ty + 3 * stat_gap_ty
stat_x0_ty = (SLIDE_WIDTH - total_w_ty) // 2
for idx, (val, label, color) in enumerate(ty_stats):
    sx = stat_x0_ty + int((stat_w_ty + stat_gap_ty) * idx)
    add_stat_card(slide, sx, Inches(3.0), stat_w_ty, stat_h_ty, val, label, color)

contact_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(4.4), SLIDE_WIDTH - Inches(1), Inches(0.5)
)
p_contact = contact_box.text_frame.paragraphs[0]
p_contact.text = "Muhammad Abdullah Tariq"
p_contact.font.size, p_contact.font.bold, p_contact.font.color.rgb, p_contact.alignment = (
    Pt(20), True, WHITE, PP_ALIGN.CENTER,
)

links_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(4.9), SLIDE_WIDTH - Inches(1), Inches(0.4)
)
p_links = links_box.text_frame.paragraphs[0]
p_links.text = "matx104.github.io/AETHERIX  |  github.com/matx104/AETHERIX"
p_links.font.size, p_links.font.color.rgb, p_links.alignment = (
    Pt(12), MED_GRAY, PP_ALIGN.CENTER,
)

bottom_bar = add_shape(
    slide, Inches(0), SLIDE_HEIGHT - Inches(0.08), SLIDE_WIDTH, Inches(0.08), shape_type=MSO_SHAPE.RECTANGLE,
    fill_color=ACCENT_BLUE,
)

# ── SLIDE ── References (Industry / Scientific) ───────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "REFERENCES \u2014 INDUSTRY & SCIENTIFIC"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.85), Inches(11), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Standards, NASA mission data, and peer-reviewed protocols"
p2.font.size, p2.font.color.rgb = Pt(14), ACCENT_CYAN

add_accent_line(slide, Inches(0.6), Inches(1.25), Inches(2.5), ACCENT_CYAN)

refs_industry_left = [
    "[1]  NASA Mars Relay Network / MRO Ka-band (0.5\u20136 Mbps)",
    "[2]  CCSDS 734.2-B-1 \u2014 Bundle Protocol Spec",
    "[3]  JPL Horizons ephemeris (Earth\u2013Mars geometry)",
    "[4]  NASA DSOC / Psyche (2023 optical demo)",
    "[5]  CCSDS 141.0-B-1 \u2014 Optical Communications",
    "[6]  IETF RFC 9171 \u2014 Bundle Protocol v7",
    "[7]  CCSDS 121.0-B-3 \u2014 Lossless Data Compression",
    "[8]  CCSDS 122.0-B-2 \u2014 Image Data Compression",
    "[9]  IETF RFC 9171 \u2014 BPv7 (primary reference)",
    "[10] IETF RFC 5326 \u2014 Licklider Transmission Protocol",
]
refs_industry_right = [
    "[11] IETF RFC 7242 \u2014 TCP Convergence Layer",
    "[12] IETF RFC 4838 \u2014 DTN Architecture",
    "[13] Bennett & Brassard, 1984 \u2014 BB84 QKD protocol",
    "[14] Ekert, 1991 \u2014 E91 entanglement-based QKD",
    "[15] Shor & Preskill, 2000 \u2014 QBER < 11% security proof",
    "[16] NIST FIPS 203 \u2014 ML-KEM (post-quantum)",
    "[17] NIST FIPS 204 \u2014 ML-DSA (post-quantum signatures)",
    "[18] BAE Systems RAD750 \u2014 radiation-hardened CPU",
    "[19] ESA / Cobham LEON3FT \u2014 fault-tolerant processor",
    "[20] NASA TMR / SEU mitigation heritage (Voyager\u2013Perseverance)",
]

left_box = add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6.3), Inches(5.0))
tf_l = left_box.text_frame
tf_l.word_wrap = True
for i, ref in enumerate(refs_industry_left):
    para = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
    para.text = ref
    para.font.size, para.font.color.rgb = Pt(11), WHITE
    para.space_after = Pt(6)

right_box = add_textbox(slide, Inches(7.0), Inches(1.5), Inches(6.3), Inches(5.0))
tf_r = right_box.text_frame
tf_r.word_wrap = True
for i, ref in enumerate(refs_industry_right):
    para = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
    para.text = ref
    para.font.size, para.font.color.rgb = Pt(11), WHITE
    para.space_after = Pt(6)

add_footer(slide, citations="Industry & scientific references [1]\u2013[20]")

# ── SLIDE ── References (Project Source Code) ─────────────────────────────────
slide = new_slide()
set_slide_bg(slide, BG_DARK)
add_slide_transition(slide, "fade")

txBox = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.6))
p = txBox.text_frame.paragraphs[0]
p.text = "REFERENCES \u2014 PROJECT SOURCE"
p.font.size, p.font.bold, p.font.color.rgb, p.alignment = Pt(28), True, WHITE, PP_ALIGN.LEFT

txBox2 = add_textbox(slide, Inches(0.6), Inches(0.85), Inches(11), Inches(0.4))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Every demonstrated number traces to a verifiable Python module"
p2.font.size, p2.font.color.rgb = Pt(14), ACCENT_ORANGE

add_accent_line(slide, Inches(0.6), Inches(1.25), Inches(2.5), ACCENT_ORANGE)

refs_proj = [
    "[A1] src/routing/rl_agent.py \u2014 Q-learning agent, \u03b5-greedy, reward function",
    "[A2] src/orbital/topology.py \u2014 5-tier network (241 nodes), BFS routing, Lagrange relays",
    "[A3] src/simulation/run_simulation.py \u2014 7-module integration test (training, delivery, rewards)",
    "[A4] src/infrastructure/link_budget.py \u2014 optical link analysis, 2\u2013200 Mbps, 1550 nm",
    "[A5] src/security/qkd.py \u2014 BB84, E91, quantum repeater, QBER detection",
    "[A6] src/computing/radiation.py \u2014 SEU/TID model, TMR, ECC, scrubbing (37,159 upsets)",
    "[A7] src/routing/prioritization.py \u2014 deadline-aware QoS scheduler, BPv7 priority classes",
    "[A8] src/simulation/run_simulation.py Module 4 \u2014 failure/recovery rewards (\u20131.438 / \u20130.201)",
]

proj_box = add_textbox(slide, Inches(0.6), Inches(1.6), Inches(12.5), Inches(4.5))
tf_p = proj_box.text_frame
tf_p.word_wrap = True
for i, ref in enumerate(refs_proj):
    para = tf_p.paragraphs[0] if i == 0 else tf_p.add_paragraph()
    para.text = ref
    para.font.size, para.font.color.rgb = Pt(13), WHITE
    para.space_after = Pt(10)

repo_box = add_textbox(slide, Inches(0.6), Inches(6.0), Inches(12.5), Inches(0.5))
p_repo = repo_box.text_frame.paragraphs[0]
p_repo.text = "All source: github.com/matx104/AETHERIX  |  Live demos: matx104.github.io/AETHERIX"
p_repo.font.size, p_repo.font.color.rgb, p_repo.alignment = Pt(12), MED_GRAY, PP_ALIGN.CENTER

add_footer(slide, citations="Project references [A1]\u2013[A8]  \u00b7  github.com/matx104/AETHERIX")

# ── SAVE BLOCK ─────────────────────────────────────────────────────────────────

for _i, _sl in enumerate(prs.slides):
    if (_i + 1) in _SPEAKER_NOTES:
        _sl.notes_slide.notes_text_frame.text = _SPEAKER_NOTES[_i + 1]

pptx_path = os.path.join(OUTPUT_DIR, "AETHERIX_Presentation.pptx")
prs.save(pptx_path)
print(f"\n\u2713 PPTX saved: {pptx_path}")
print(f"  Slides: {len(prs.slides)}")
